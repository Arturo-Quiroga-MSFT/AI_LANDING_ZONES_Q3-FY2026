"""
Generate an editable PPTX slide deck for Workshop 3: Landing Zones to Production (GenAIOps Bridge).
Uses python-pptx to create real text, tables, and callouts.
Output: SLIDE-DECK-EDITABLE.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Microsoft brand colors
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x5A, 0x9E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARNING_ORANGE = RGBColor(0xD8, 0x3B, 0x01)
WARNING_BG = RGBColor(0xFF, 0xF4, 0xCE)
GREEN = RGBColor(0x10, 0x7C, 0x10)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_footer(slide, text="AI Center of Excellence V2 | Q3 FY2026"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY


def add_blue_bar(slide, top=Inches(1.35), width=Inches(2.5)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), top, width, Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZURE_BLUE
    bar.line.fill.background()


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def make_title_slide(title, subtitle="", footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, AZURE_BLUE)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    if footer_text:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
        tf2 = txBox2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = footer_text
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
        p3.alignment = PP_ALIGN.CENTER

    return slide


def make_section_slide(module_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Module {module_num}" if module_num else ""
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    return slide


def make_content_slide(title, bullets=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()

            if " \u2014 " in bullet:
                parts = bullet.split(" \u2014 ", 1)
                run1 = p2.add_run()
                run1.text = parts[0] + " \u2014 "
                run1.font.bold = True
                run1.font.size = Pt(18)
                run1.font.color.rgb = DARK_GRAY
                run2 = p2.add_run()
                run2.text = parts[1]
                run2.font.size = Pt(18)
                run2.font.color.rgb = DARK_GRAY
            else:
                p2.text = bullet
                p2.font.size = Pt(18)
                p2.font.color.rgb = DARK_GRAY

            p2.space_before = Pt(8)
            p2.level = 0

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_two_column_slide(title, left_title, left_bullets, right_title, right_bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(22)
    lp.font.color.rgb = DARK_BLUE
    lp.font.bold = True
    for b in left_bullets:
        lp2 = ltf.add_paragraph()
        lp2.text = b
        lp2.font.size = Pt(16)
        lp2.font.color.rgb = DARK_GRAY
        lp2.space_before = Pt(6)

    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(22)
    rp.font.color.rgb = DARK_BLUE
    rp.font.bold = True
    for b in right_bullets:
        rp2 = rtf.add_paragraph()
        rp2.text = b
        rp2.font.size = Pt(16)
        rp2.font.color.rgb = DARK_GRAY
        rp2.space_before = Pt(6)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_table_slide(title, headers, rows, notes="", col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(11.9)
    table_height = Inches(min(num_rows * 0.5, 5.0))

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.7), Inches(1.6),
        table_width, table_height
    )
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZURE_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = DARK_GRAY

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_callout_slide(title, callout_text, bullets=None, callout_color="blue", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    bar_color = AZURE_BLUE if callout_color == "blue" else WARNING_ORANGE
    bg_color = RGBColor(0xE8, 0xF4, 0xFD) if callout_color == "blue" else WARNING_BG

    callout_top = Inches(1.6)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), callout_top, Inches(11.9), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), callout_top, Inches(0.06), Inches(1.2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = bar_color
    accent.line.fill.background()

    ctxBox = slide.shapes.add_textbox(
        Inches(1.0), callout_top + Inches(0.15), Inches(11.3), Inches(0.9)
    )
    ctf = ctxBox.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = callout_text
    cp.font.size = Pt(16)
    cp.font.color.rgb = DARK_GRAY
    cp.font.italic = True

    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(11.9), Inches(4.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(18)
            p2.font.color.rgb = DARK_GRAY
            p2.space_before = Pt(8)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_code_slide(title, code_text, bullets_after=None, notes=""):
    """Create a slide with a code block and optional bullets below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    code_top = Inches(1.6)
    code_height = Inches(2.4)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), code_top, Inches(11.9), code_height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    box.line.fill.background()

    code_box = slide.shapes.add_textbox(
        Inches(1.0), code_top + Inches(0.15), Inches(11.3), code_height - Inches(0.3)
    )
    ctf = code_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = code_text
    cp.font.size = Pt(14)
    cp.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
    cp.font.name = "Consolas"

    if bullets_after:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets_after):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(16)
            p2.font.color.rgb = DARK_GRAY
            p2.space_before = Pt(6)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# ============================================================
# SLIDE GENERATION
# ============================================================

# ---- SECTION 0: WELCOME & CONTEXT ----

# --- SLIDE 1: Title ---
make_title_slide(
    "Landing Zones to Production",
    "The GenAIOps Bridge \u2014 Partner Enablement Workshop (Advanced)",
    "AI Center of Excellence V2 | Q3 FY2026"
)

# --- SLIDE 2: Agenda ---
make_table_slide(
    "Agenda",
    ["Time", "Module", "Description"],
    [
        ["0:00", "Welcome & Context", "Workshop 2 recap, production readiness gap"],
        ["0:15", "Module 1", "CI/CD for AI Workloads"],
        ["1:05", "\u2615 Break", "10 minutes"],
        ["1:15", "Module 2", "Production Monitoring & Observability"],
        ["2:00", "\u2615 Break", "10 minutes"],
        ["2:10", "Module 3", "Agent Governance & Security"],
        ["2:55", "\u2615 Break", "10 minutes"],
        ["3:05", "Module 4", "Cost Optimization & Scaling"],
        ["3:35", "Module 5", "Wrap-up & Production Planning"],
    ],
    notes="Walk through the agenda. This is the most architecture-heavy workshop in the series.",
    col_widths=[1.5, 3.0, 7.4]
)

# --- SLIDE 3: Where We Are ---
make_content_slide(
    "Where We Are in the Journey",
    [
        "WS1: UNDERSTAND \u2014 Concepts, architectures, decision frameworks",
        "WS2: BUILD \u2014 Deploy LZ, RAG app, explore agents",
        "WS3: OPERATE \u2190 you are here \u2014 CI/CD, monitoring, governance, cost",
        "",
        "You have a deployed workload. Now: how do you run it in production?",
        "",
        "Deploying is Day 1. Operating is Day 2 through Day N.",
    ],
    notes="Frame the workshop: WS2 got them deployed, WS3 gets them production-ready."
)

# ---- SECTION 1: CI/CD FOR AI WORKLOADS ----

# --- SLIDE 4: Module 1 divider ---
make_section_slide("1", "CI/CD for AI Workloads")

# --- SLIDE 5: The Two Pipelines ---
make_two_column_slide(
    "The Two Pipelines",
    "Infrastructure Pipeline",
    [
        "\u2022 Deploys/updates Landing Zone",
        "\u2022 Bicep or Terraform",
        "\u2022 Trigger: PR merge to infra/",
        "\u2022 Changes: networking, identity, PaaS config",
        "\u2022 Slower cadence, higher risk",
    ],
    "Application Pipeline",
    [
        "\u2022 Builds/deploys AI app",
        "\u2022 Container build + model config",
        "\u2022 Trigger: PR merge to src/",
        "\u2022 Changes: app code, prompts, model versions",
        "\u2022 Faster cadence, more frequent",
    ],
    notes="These are separate pipelines with different triggers and approval gates."
)

# --- SLIDE 6: Infrastructure Pipeline ---
make_code_slide(
    "Infrastructure Pipeline Pattern",
    "# Trigger: PR merge to infra/\nsteps:\n  - Lint (bicep lint / terraform validate)\n  - Validate (bicep build / terraform plan)\n  - What-if / Plan (preview changes)\n  - Approve (manual gate)\n  - Deploy (bicep deploy / terraform apply)",
    [
        "\u2022 Preview all changes before applying (what-if / plan)",
        "\u2022 Manual approval gate before production",
        "\u2022 Immutable deployments \u2014 no manual portal changes",
        "\u2022 IaC drift detection on schedule",
    ],
    notes="Show a real GitHub Actions or Azure DevOps pipeline YAML if available."
)

# --- SLIDE 7: Application Pipeline ---
make_code_slide(
    "Application Pipeline Pattern",
    "# Trigger: PR merge to src/\nsteps:\n  - Build container image\n  - Run unit tests\n  - Run model evaluation (groundedness, relevance)\n  - Push to Azure Container Registry\n  - Deploy to Container Apps (staging)\n  - Integration tests\n  - Approve (manual gate)\n  - Deploy to Container Apps (production)",
    [
        "Key addition: Model evaluation step BEFORE production deploy.",
    ],
    notes="The model evaluation step is what makes this different from a standard CI/CD pipeline."
)

# --- SLIDE 8: Model Deployment & Versioning ---
make_table_slide(
    "Model Deployment & Versioning",
    ["What to Version", "Strategy"],
    [
        ["Model", "Track model name + version (e.g., gpt-4o-2024-08-06)"],
        ["System prompts", "Store in repo as code, version with git"],
        ["RAG configuration", "Index schema, chunking strategy, embedding model"],
        ["Parameters", "Temperature, top-p, max tokens as config files"],
    ],
    notes="Deployment patterns: blue/green for model swaps, canary (5% traffic), always keep previous version deployable.",
    col_widths=[3.0, 8.9]
)

# --- SLIDE 9: Testing AI Workloads ---
make_table_slide(
    "Testing AI Workloads",
    ["Test Level", "What It Tests", "Tools"],
    [
        ["Unit tests", "Deterministic logic, API contracts", "pytest, Jest"],
        ["Integration tests", "End-to-end: query \u2192 search \u2192 model \u2192 response", "Custom test harness"],
        ["Model evaluation", "Response quality, groundedness, relevance", "AI Foundry evaluation"],
        ["Load testing", "Token throughput, latency under concurrency", "Azure Load Testing, k6"],
    ],
    notes="Model evaluation is not optional. A code change that passes unit tests can still degrade response quality.",
    col_widths=[2.5, 5.0, 4.4]
)

# --- SLIDE 10: Environments & Promotion ---
make_content_slide(
    "Environments & Promotion",
    [
        "Dev  \u2192  Staging  \u2192  Production",
        "",
        "\u2022 Dev: Experiment, iterate, lower SKUs, auto-deploy",
        "\u2022 Staging: Validate, load test, model eval, auto + gate",
        "\u2022 Production: Serve, monitor, full SKUs, gated only",
        "",
        "Key practices:",
        "\u2022 Same IaC templates across all environments, different parameter files",
        "\u2022 Approval gates between staging \u2192 production",
        "\u2022 Feature flags for gradual rollout of new agent capabilities",
        "\u2022 Environment-specific: model SKU, search tier, network config",
    ],
    notes="Never skip staging. Model evaluation in staging catches issues that unit tests miss."
)

# --- SLIDE 11: Key Takeaways Module 1 ---
make_content_slide(
    "Key Takeaways \u2014 Module 1",
    [
        "\u2705 Separate infrastructure and application pipelines",
        "\u2705 Treat prompts and model config as versioned code",
        "\u2705 Test AI workloads at multiple levels (unit \u2192 integration \u2192 model eval \u2192 load)",
        "\u2705 Automate deployments but gate production promotions",
        "\u2705 Blue/green or canary for model swaps \u2014 never yolo to production",
    ],
    notes="Transition to break."
)

# --- SLIDE 12: Break ---
make_title_slide("\u2615 Break", "10 minutes")

# ---- SECTION 2: PRODUCTION MONITORING & OBSERVABILITY ----

# --- SLIDE 13: Module 2 divider ---
make_section_slide("2", "Production Monitoring & Observability")

# --- SLIDE 14: Monitoring Stack Overview ---
make_table_slide(
    "Monitoring Stack Overview",
    ["Tool", "What It Covers"],
    [
        ["Application Insights", "Request tracing, errors, latency, dependencies"],
        ["Log Analytics", "Centralized logs, KQL queries, diagnostic settings"],
        ["AI Foundry Tracing", "Model-level call traces, token usage, response quality"],
        ["Azure Workbooks / Grafana", "Custom dashboards"],
        ["Azure Monitor Alerts", "Proactive notification"],
    ],
    notes="All already deployed from azd up in Workshop 2 \u2014 now we configure them for production.",
    col_widths=[4.0, 7.9]
)

# --- SLIDE 15: AI Monitoring Pyramid ---
make_content_slide(
    "The AI Monitoring Pyramid",
    [
        "                        Business KPIs",
        "                     /                  \\",
        "              Model Quality        User Experience",
        "           /                                        \\",
        "   Infrastructure Health              Cost Metrics",
        "",
        "Monitor bottom-up: Infrastructure health \u2192 Model quality \u2192 Business value",
        "",
        "If infrastructure is degraded, model metrics are meaningless.",
        "If model quality drops, business KPIs will follow.",
    ],
    notes="This pyramid is the organizing principle for all monitoring decisions."
)

# --- SLIDE 16: Key Metrics Dashboard ---
make_table_slide(
    "Key Metrics Dashboard",
    ["Category", "Metric", "Target"],
    [
        ["Latency", "Response P50 / P95 / P99", "< 2s / < 5s / < 10s"],
        ["Tokens", "Prompt + completion per request", "Monitor trend"],
        ["Cost", "Cost per 1K requests", "< budget threshold"],
        ["Errors", "Error rate by category", "< 1%"],
        ["Retrieval", "% relevant search results", "> 80%"],
        ["Availability", "Uptime %", "> 99.5%"],
    ],
    notes="Build a single dashboard that shows all of these. If you can't see it at a glance, you won't catch degradation early.",
    col_widths=[2.5, 5.0, 4.4]
)

# --- SLIDE 17: Agent-Specific Observability ---
make_content_slide(
    "Agent-Specific Observability",
    [
        "When running AI agents (vs. classic RAG), monitor additional dimensions:",
        "",
        "\u2022 Tool call monitoring \u2014 which tools called, success/failure rate, latency per tool",
        "\u2022 Reasoning chain tracing \u2014 full decision path per request",
        "\u2022 Escalation tracking \u2014 when agent hands off to human, frequency, reasons",
        "\u2022 Drift detection \u2014 are agent behaviors changing over time?",
        "",
        "\U0001f4da CAF Manage AI Agents: learn.microsoft.com/azure/cloud-adoption-framework/ai-agents/manage",
    ],
    notes="These are the metrics that differentiate monitoring agents from monitoring RAG."
)

# --- SLIDE 18: Alerting Rules ---
make_table_slide(
    "Alerting Rules",
    ["Alert", "Threshold", "Severity"],
    [
        ["Response latency P95", "> 5000 ms", "\u26a0\ufe0f Warning"],
        ["Error rate", "> 2% sustained 5 min", "\U0001f534 Critical"],
        ["Token budget", "> 80% of daily limit", "\u26a0\ufe0f Warning"],
        ["Model drift score", "> threshold", "\u26a0\ufe0f Warning"],
        ["Availability", "< 99.5% in 1 hour", "\U0001f534 Critical"],
        ["Hallucination rate", "> 5% of responses", "\U0001f534 Critical"],
    ],
    notes="These are starting points. Tune thresholds after 2 weeks of production data.",
    col_widths=[4.0, 4.0, 3.9]
)

# --- SLIDE 19: SLA Definition ---
make_table_slide(
    "SLA Definition for AI Workloads",
    ["SLA Type", "Metric", "Example Target"],
    [
        ["Availability", "Uptime %", "99.5%"],
        ["Latency", "P95 response time", "< 5 seconds"],
        ["Quality", "% grounded responses", "> 95%"],
        ["Cost", "Max cost per 1K requests", "< $X"],
        ["Agent", "% tasks completed without escalation", "> 85%"],
    ],
    notes="Partner talking point: Your AI SLA should include response quality, not just uptime.",
    col_widths=[2.5, 5.0, 4.4]
)

# --- SLIDE 20: Key Takeaways Module 2 ---
make_content_slide(
    "Key Takeaways \u2014 Module 2",
    [
        "\u2705 AI monitoring = traditional monitoring + model-specific metrics",
        "\u2705 Build dashboards around the monitoring pyramid (infra \u2192 model \u2192 business)",
        "\u2705 Define SLAs that include quality, not just availability",
        "\u2705 Agent workloads need tool call and reasoning observability",
        "\u2705 Alerting thresholds are starting points \u2014 tune with production data",
    ],
    notes="Transition to break."
)

# --- SLIDE 21: Break ---
make_title_slide("\u2615 Break", "10 minutes")

# ---- SECTION 3: AGENT GOVERNANCE & SECURITY ----

# --- SLIDE 22: Module 3 divider ---
make_section_slide("3", "Agent Governance & Security")

# --- SLIDE 23: CAF Govern & Secure Phase ---
make_content_slide(
    "CAF Govern & Secure Phase",
    [
        "The 4-phase CAF AI Agent lifecycle:",
        "",
        "Plan  \u2192  Govern & Secure  \u2192  Build  \u2192  Manage",
        "                  \u2191",
        "            YOU ARE HERE",
        "",
        "This phase happens BEFORE you build agents in production.",
        "Governance is not an afterthought \u2014 it's a prerequisite.",
        "",
        "\U0001f4da CAF Govern & Secure: learn.microsoft.com/azure/cloud-adoption-framework/ai-agents/govern-secure",
    ],
    notes="Emphasize: governance decisions made now determine what's possible later."
)

# --- SLIDE 24: Data Access Controls ---
make_table_slide(
    "Data Access Controls",
    ["Control", "Implementation"],
    [
        ["Document-level", "AI Search security filters"],
        ["Row-level", "Cosmos DB resource tokens"],
        ["Service-level", "RBAC per managed identity"],
        ["Network-level", "Private endpoints, NSGs"],
    ],
    notes="Apply principle of least privilege to agent data access. Classify data sources: which are agent-accessible, which are not.",
    col_widths=[3.5, 8.4]
)

# --- SLIDE 25: Tool Permission Boundaries ---
make_table_slide(
    "Tool Permission Boundaries",
    ["Agent Type", "Allowed Tools", "Restrictions"],
    [
        ["Productivity", "AI Search, knowledge base", "Read-only, no external calls"],
        ["Action", "Specific APIs, ticketing, CRM", "Allowlisted endpoints, rate-limited"],
        ["Automation", "Full orchestration suite", "Budget caps, approval for sensitive ops"],
    ],
    notes="Default deny. Tool allowlists per agent type. Rate limiting and budget caps per tool.",
    col_widths=[2.5, 4.5, 4.9]
)

# --- SLIDE 26: Human-in-the-Loop Policies ---
make_callout_slide(
    "Human-in-the-Loop Policies",
    "Human-in-the-loop policies must be defined BEFORE agents are deployed to production. Retrofitting is expensive and risky.",
    [
        "When must agents escalate to humans?",
        "",
        "\u2022 Confidence threshold \u2014 agent is uncertain about intent or correct action",
        "\u2022 Sensitive actions \u2014 financial transactions, PII processing, data deletion",
        "\u2022 Compliance triggers \u2014 regulated actions require human approval",
        "\u2022 Accumulated risk \u2014 multiple borderline decisions in one session",
        "",
        "Audit and log every escalation decision \u2014 both escalated and not-escalated.",
    ],
    callout_color="orange",
    notes="This is a design decision, not an afterthought."
)

# --- SLIDE 27: Responsible AI Guardrails ---
make_table_slide(
    "Responsible AI Guardrails",
    ["Guardrail", "Purpose", "Tool"],
    [
        ["Content filtering", "Block unsafe outputs", "Azure AI Content Safety"],
        ["Grounding enforcement", "Agents must cite sources", "System prompt + evaluation"],
        ["Hallucination detection", "Flag ungrounded responses", "AI Foundry evaluation"],
        ["Bias monitoring", "Detect unfair decision patterns", "Custom metrics + dashboards"],
    ],
    notes="Design Checklist alignment: S-R4 (MITRE ATLAS + OWASP), S-R5 (AI Content Safety), G-R3 (Responsible AI dashboard).",
    col_widths=[3.0, 4.5, 4.4]
)

# --- SLIDE 28: Multi-Agent Security ---
make_two_column_slide(
    "Multi-Agent Security",
    "Must Have",
    [
        "\u2022 Credential isolation between agents",
        "\u2022 Per-agent managed identities",
        "\u2022 Separate RBAC scopes",
        "\u2022 Network segmentation",
    ],
    "Must Design",
    [
        "\u2022 Cross-boundary trust models",
        "\u2022 State synchronization security",
        "\u2022 Inter-agent communication auth",
        "\u2022 Shared resource access patterns",
    ],
    notes="When NOT to use multi-agent: if complexity doesn't justify value. Most workloads start fine with single agent."
)

# --- SLIDE 29: Agent Audit Trail ---
make_content_slide(
    "Agent Audit Trail",
    [
        "Log everything the agent does:",
        "",
        "\u2022 Every tool call (input, output, latency)",
        "\u2022 Every data access (what was queried, what was returned)",
        "\u2022 Every decision (why this tool, why this response)",
        "\u2022 Every escalation (to whom, why, outcome)",
        "\u2022 Every output (what was sent to the user)",
        "",
        "Retention: Follow compliance requirements (min 90 days recommended)",
        "Forensics: Must be able to reconstruct any agent decision chain after the fact",
    ],
    notes="In regulated industries, this is not optional. Even for non-regulated, it's critical for debugging."
)

# --- SLIDE 30: Key Takeaways Module 3 ---
make_content_slide(
    "Key Takeaways \u2014 Module 3",
    [
        "\u2705 Governance before deployment, not after",
        "\u2705 Least privilege for data access and tool permissions",
        "\u2705 Human-in-the-loop is a design decision, not an afterthought",
        "\u2705 Responsible AI guardrails are mandatory, not optional",
        "\u2705 Audit everything \u2014 tool calls, decisions, escalations, outputs",
        "",
        "Exercise: Walk through Design Checklist governance (G-R1 to G-R5)",
        "and security (S-R1 to S-R5) against your deployed environment.",
    ],
    notes="If time allows, have participants do the exercise in pairs."
)

# --- SLIDE 31: Break ---
make_title_slide("\u2615 Break", "10 minutes")

# ---- SECTION 4: COST OPTIMIZATION & SCALING ----

# --- SLIDE 32: Module 4 divider ---
make_section_slide("4", "Cost Optimization & Scaling")

# --- SLIDE 33: AI Workload Cost Model ---
make_table_slide(
    "AI Workload Cost Model",
    ["Cost Component", "Driver", "Optimization Lever"],
    [
        ["Model inference", "Tokens (prompt + completion)", "Smaller models, shorter prompts, caching"],
        ["AI Search", "Tier + query volume", "Right-size tier, optimize index"],
        ["Compute", "Container App instances", "Autoscaling, right-size SKU"],
        ["Storage", "Document + chat history", "Lifecycle policies, archive old data"],
        ["Networking", "Private endpoints, egress", "Minimize cross-region traffic"],
    ],
    notes="Official reference: AI Landing Zones Cost Guide and Azure Pricing Calculator estimate.",
    col_widths=[2.5, 4.5, 4.9]
)

# --- SLIDE 34: Model Selection Economics ---
make_table_slide(
    "Model Selection Economics",
    ["Model", "Relative Cost", "Best For"],
    [
        ["GPT-4o", "$$$", "Complex reasoning, multi-step tasks"],
        ["GPT-4o-mini", "$$", "Most production workloads, good cost/quality"],
        ["GPT-3.5-turbo", "$", "Simple queries, high volume, low latency"],
    ],
    notes="Model selection is the single biggest cost lever. Route simple queries to smaller models.",
    col_widths=[3.0, 3.0, 5.9]
)

# --- SLIDE 35: Model Selection Strategies ---
make_content_slide(
    "Model Cost Optimization Strategies",
    [
        "\u2022 Route \u2014 simple queries to smaller models, complex to larger",
        "\u2022 Cache \u2014 responses for common/repeated queries",
        "\u2022 Optimize prompts \u2014 reduce token count without losing quality",
        "\u2022 Right-size context \u2014 don't stuff the full document when a summary suffices",
        "",
        "Model selection is the single biggest cost lever for AI workloads.",
    ],
    notes="Show concrete examples of prompt optimization: before/after token counts."
)

# --- SLIDE 36: Scaling Patterns ---
make_table_slide(
    "Scaling Patterns",
    ["Component", "Scaling Mechanism", "Key Config"],
    [
        ["Container Apps", "HTTP-based, KEDA, custom", "Min/max replicas, scale rules"],
        ["Azure OpenAI", "PTU (reserved) vs. PAYGO", "PTU for baseline, PAYGO for bursts"],
        ["AI Search", "Replica scaling", "Read replicas for query throughput"],
        ["Cosmos DB", "Autoscale RUs", "Min/max RUs, partition key design"],
    ],
    notes="PTU vs. PAYGO: Use Provisioned Throughput Units for predictable baseline traffic. Use PAYGO for spikes. Combine both.",
    col_widths=[2.5, 4.5, 4.9]
)

# --- SLIDE 37: Cost Allocation & Showback ---
make_table_slide(
    "Cost Allocation & Showback",
    ["Tag", "Purpose", "Example"],
    [
        ["team", "Cost allocation", "team:ai-platform"],
        ["project", "Project tracking", "project:customer-support-agent"],
        ["environment", "Env separation", "environment:production"],
        ["workload", "Workload type", "workload:rag-chat"],
    ],
    notes="Set up Azure Cost Management budgets per tag group. Track per-request cost. Review weekly first month, monthly after.",
    col_widths=[2.0, 4.0, 5.9]
)

# --- SLIDE 38: Key Takeaways Module 4 ---
make_content_slide(
    "Key Takeaways \u2014 Module 4",
    [
        "\u2705 Understand the cost model before deploying to production",
        "\u2705 Model selection is the biggest cost lever",
        "\u2705 Autoscale everything, right-size everything",
        "\u2705 Combine PTU + PAYGO for optimal model inference cost",
        "\u2705 Tag resources for showback from day one",
    ],
    notes="Transition to Module 5: Wrap-up."
)

# ---- SECTION 5: WRAP-UP & PRODUCTION PLANNING ----

# --- SLIDE 39: Module 5 divider ---
make_section_slide("5", "Wrap-up & Production Planning")

# --- SLIDE 40: Production Readiness Checklist ---
make_table_slide(
    "Production Readiness Checklist",
    ["Area", "Item", "Status"],
    [
        ["CI/CD", "Infrastructure pipeline (Bicep/Terraform)", "\u2610"],
        ["CI/CD", "Application pipeline (build, test, deploy)", "\u2610"],
        ["CI/CD", "Model evaluation in pipeline", "\u2610"],
        ["Monitoring", "Dashboards configured (metrics pyramid)", "\u2610"],
        ["Monitoring", "Alerting rules active", "\u2610"],
        ["Governance", "Agent data access policies defined", "\u2610"],
        ["Governance", "Tool permission boundaries enforced", "\u2610"],
        ["Governance", "Human-in-the-loop policies configured", "\u2610"],
        ["Governance", "Responsible AI review completed", "\u2610"],
        ["Cost", "Budgets and scaling limits set", "\u2610"],
        ["Operations", "Operational runbooks documented", "\u2610"],
        ["Operations", "Incident response plan defined", "\u2610"],
        ["Security", "Load testing passed", "\u2610"],
        ["Security", "Security review completed", "\u2610"],
    ],
    notes="Use this checklist as a gate before going to production.",
    col_widths=[2.0, 8.0, 1.9]
)

# --- SLIDE 41: Operational Runbook Template ---
make_table_slide(
    "Operational Runbook Template",
    ["Element", "Details"],
    [
        ["Incident classification", "P1-P4 for AI workloads"],
        ["AI failure modes", "Tool failures, model degradation, hallucination spike"],
        ["Escalation paths", "On-call rotation, subject matter experts"],
        ["Recovery procedures", "Per failure type: rollback model, restart service, failover"],
        ["Post-incident review", "Root cause, impact, remediation, prevention"],
    ],
    notes="Agent-specific failures: tool invocation failures, reasoning loops, confidence collapse, escalation floods.",
    col_widths=[3.5, 8.4]
)

# --- SLIDE 42: Workshop Series Recap ---
make_table_slide(
    "The Full Journey \u2014 Workshop Series Recap",
    ["Workshop", "Theme", "Outcome"],
    [
        ["WS1", "Understand", "Architecture decisions, design checklist, deployment paths"],
        ["WS2", "Build", "Deployed Landing Zone, RAG app, agent exploration"],
        ["WS3", "Operate", "CI/CD, monitoring, governance, cost optimization"],
    ],
    notes="Partner outcome: end-to-end capability to take customers from concept to production.",
    col_widths=[2.0, 2.5, 7.4]
)

# --- SLIDE 43: Resources ---
make_table_slide(
    "Key Resources",
    ["Resource", "Link"],
    [
        ["AI Landing Zones Repo", "github.com/Azure/AI-Landing-Zones"],
        ["Design Checklist", "AI-Landing-Zones-Design-Checklist.md"],
        ["Cost Guide", "AI-Landing-Zones-Cost-Guide.md"],
        ["CAF AI Agent Adoption", "learn.microsoft.com/azure/cloud-adoption-framework/ai-agents/"],
        ["CAF Govern & Secure", "learn.microsoft.com/.../ai-agents/govern-secure"],
        ["CAF Manage Agents", "learn.microsoft.com/.../ai-agents/manage"],
        ["Azure Monitor Baseline Alerts", "azure.github.io/azure-monitor-baseline-alerts/"],
    ],
    col_widths=[4.0, 7.9]
)

# --- SLIDE 44: Q&A ---
make_title_slide(
    "Q&A",
    "Questions? Feedback?",
    "AI Center of Excellence V2 | Partner Enablement Team"
)

# --- SLIDE 45: Thank You ---
make_title_slide(
    "Thank You",
    "Workshop materials available in the partner repository\nFeedback: Survey link in chat\nContact: Arturo Quiroga (PSA) | Anahita Afshari (PSA)",
    "AI Center of Excellence V2 | Q3 FY2026"
)


# ============================================================
# SAVE
# ============================================================

output_path = os.path.join(SCRIPT_DIR, "SLIDE-DECK-EDITABLE.pptx")
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to {output_path}")
