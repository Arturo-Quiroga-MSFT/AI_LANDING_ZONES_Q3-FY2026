"""
Generate an editable PPTX slide deck for Workshop 2: Deploying Your First Gen AI Workload.
Uses python-pptx to create real text, tables, and callouts.
Output: SLIDE-DECK-EDITABLE.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Microsoft brand colors
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x5A, 0x9E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARNING_ORANGE = RGBColor(0xD8, 0x3B, 0x01)
WARNING_BG = RGBColor(0xFF, 0xF4, 0xCE)
GREEN = RGBColor(0x10, 0x7C, 0x10)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_footer(slide, text="AI Center of Excellence V2 | Q3 FY2026"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY


def add_blue_bar(slide, top=Inches(1.35), width=Inches(2.5)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), top, width, Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZURE_BLUE
    bar.line.fill.background()


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def make_title_slide(title, subtitle="", footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, AZURE_BLUE)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    if footer_text:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
        tf2 = txBox2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = footer_text
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
        p3.alignment = PP_ALIGN.CENTER

    return slide


def make_section_slide(module_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Module {module_num}" if module_num else ""
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    return slide


def make_content_slide(title, bullets=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()

            if " \u2014 " in bullet:
                parts = bullet.split(" \u2014 ", 1)
                run1 = p2.add_run()
                run1.text = parts[0] + " \u2014 "
                run1.font.bold = True
                run1.font.size = Pt(18)
                run1.font.color.rgb = DARK_GRAY
                run2 = p2.add_run()
                run2.text = parts[1]
                run2.font.size = Pt(18)
                run2.font.color.rgb = DARK_GRAY
            else:
                p2.text = bullet
                p2.font.size = Pt(18)
                p2.font.color.rgb = DARK_GRAY

            p2.space_before = Pt(8)
            p2.level = 0

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_two_column_slide(title, left_title, left_bullets, right_title, right_bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(22)
    lp.font.color.rgb = DARK_BLUE
    lp.font.bold = True
    for b in left_bullets:
        lp2 = ltf.add_paragraph()
        lp2.text = b
        lp2.font.size = Pt(16)
        lp2.font.color.rgb = DARK_GRAY
        lp2.space_before = Pt(6)

    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(22)
    rp.font.color.rgb = DARK_BLUE
    rp.font.bold = True
    for b in right_bullets:
        rp2 = rtf.add_paragraph()
        rp2.text = b
        rp2.font.size = Pt(16)
        rp2.font.color.rgb = DARK_GRAY
        rp2.space_before = Pt(6)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_table_slide(title, headers, rows, notes="", col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(11.9)
    table_height = Inches(min(num_rows * 0.5, 5.0))

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.7), Inches(1.6),
        table_width, table_height
    )
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZURE_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = DARK_GRAY

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_callout_slide(title, callout_text, bullets=None, callout_color="blue", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    bar_color = AZURE_BLUE if callout_color == "blue" else WARNING_ORANGE
    bg_color = RGBColor(0xE8, 0xF4, 0xFD) if callout_color == "blue" else WARNING_BG

    callout_top = Inches(1.6)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), callout_top, Inches(11.9), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), callout_top, Inches(0.06), Inches(1.2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = bar_color
    accent.line.fill.background()

    ctxBox = slide.shapes.add_textbox(
        Inches(1.0), callout_top + Inches(0.15), Inches(11.3), Inches(0.9)
    )
    ctf = ctxBox.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = callout_text
    cp.font.size = Pt(16)
    cp.font.color.rgb = DARK_GRAY
    cp.font.italic = True

    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(11.9), Inches(4.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(18)
            p2.font.color.rgb = DARK_GRAY
            p2.space_before = Pt(8)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_code_slide(title, code_text, bullets_after=None, notes=""):
    """Create a slide with a code block and optional bullets below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    # Code box background
    code_top = Inches(1.6)
    code_height = Inches(2.4)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), code_top, Inches(11.9), code_height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    box.line.fill.background()

    code_box = slide.shapes.add_textbox(
        Inches(1.0), code_top + Inches(0.15), Inches(11.3), code_height - Inches(0.3)
    )
    ctf = code_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = code_text
    cp.font.size = Pt(14)
    cp.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
    cp.font.name = "Consolas"

    if bullets_after:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets_after):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(16)
            p2.font.color.rgb = DARK_GRAY
            p2.space_before = Pt(6)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# ============================================================
# SLIDE GENERATION
# ============================================================

# --- SLIDE 1: Title ---
make_title_slide(
    "From RAG to Agents",
    "Deploying Your First Gen AI Workload \u2014 Partner Enablement Workshop (Hands-On)",
    "AI Center of Excellence V2 | Q3 FY2026"
)

# --- SLIDE 2: Agenda ---
make_table_slide(
    "Agenda",
    ["Time", "Module", "Description"],
    [
        ["0:00", "Welcome & Setup", "Verify prerequisites, agenda"],
        ["0:15", "Module 1", "Quick Recap & Deployment Strategy"],
        ["0:35", "Module 2 \U0001f52c", "Deploy the Landing Zone (HANDS-ON)"],
        ["1:35", "\u2615 Break", "15 minutes"],
        ["1:50", "Module 3 \U0001f52c", "Configure & Deploy RAG App (HANDS-ON)"],
        ["2:30", "Module 4", "From RAG to Agents"],
        ["3:00", "Module 5", "Monitoring & Observability"],
        ["3:20", "Module 6", "Wrap-up & Next Steps"],
    ],
    notes="Walk through the agenda. Highlight the two hands-on lab modules.",
    col_widths=[1.5, 3.0, 7.4]
)

# --- SLIDE 3: Learning Objectives ---
make_content_slide(
    "Learning Objectives",
    [
        "Deploy a complete AI Landing Zone using azd up",
        "Configure AI Foundry with standard (private) setup",
        "Deploy a sample RAG chat application end-to-end",
        "Distinguish RAG vs. AI agents in practice",
        "Explore Microsoft Foundry agent capabilities",
        "Implement basic monitoring and observability",
        "Assess when a customer should graduate from RAG to agents",
    ],
    notes="These objectives are measurable. We'll validate each one during the workshop."
)

# --- SLIDE 4: Setup Validation ---
make_callout_slide(
    "Setup Validation Checklist",
    "\u26a0\ufe0f Run these now \u2014 raise your hand if anything fails:",
    [
        "az --version          # Need 2.61.0+",
        "azd version           # Need 1.15.0+",
        "az account show       # Must be logged in",
        "az cognitiveservices usage list --location eastus2 -o table",
        "",
        "\u2610 Azure CLI \u2265 2.61.0",
        "\u2610 Azure Developer CLI \u2265 1.15.0",
        "\u2610 Logged in to correct subscription",
        "\u2610 Sufficient OpenAI quota in target region",
        "\u2610 Owner or Contributor + UAA role",
    ],
    callout_color="orange",
    notes="If anything fails: pair with a neighbor, or follow along with the instructor demo."
)

# --- SLIDE 5: Module 1 divider ---
make_section_slide("1", "Quick Recap & Deployment Strategy")

# --- SLIDE 6: Workshop 1 Recap ---
make_content_slide(
    "Workshop 1 Recap (60-Second Version)",
    [
        "AI Landing Zone \u2014 production-ready application landing zone for AI workloads",
        "Two architectures \u2014 with Platform LZ (enterprise) / without (standalone)",
        "Design Checklist \u2014 40+ recommendations across 10 design areas",
        "Four deployment paths \u2014 azd, Bicep, Terraform, Portal (Deploy to Azure)",
        "",
        "Today: We deploy and build on top of this foundation",
        "",
        "\U0001f4da Design Checklist | Cost Guide \u2014 see links in chat",
    ],
    notes="Quick refresher. Don't spend more than 2 minutes here."
)

# --- SLIDE 7: Today's Narrative ---
make_content_slide(
    "Today's Narrative \u2014 RAG to Agents",
    [
        "Deploy Infrastructure  \u2192  Build Baseline Workload  \u2192  Explore What's Next  \u2192  Operate",
        "",
        "  AI Landing Zone       Classic RAG App          AI Agents            Monitoring",
        "  (azd up)              (chat + search)          (when ready)         & Ops",
        "",
        "The key question partners must answer:",
        '"Does this customer need agents, or is RAG enough?"',
    ],
    notes="Frame the entire workshop around this narrative arc."
)

# --- SLIDE 8: What Gets Deployed ---
make_table_slide(
    "What Gets Deployed (~30+ Resources)",
    ["Layer", "Resources"],
    [
        ["AI Services", "AI Foundry (hub + project), Azure OpenAI, AI Search"],
        ["Data", "Cosmos DB (chat history), Storage Account (documents)"],
        ["Compute", "Container Apps (app runtime)"],
        ["Security", "Key Vault, Managed Identities, Private Endpoints, NSGs"],
        ["Networking", "Virtual Network, Subnets, Private DNS Zones"],
        ["Monitoring", "Log Analytics, Application Insights"],
        ["Governance", "Microsoft Fabric, Microsoft Purview"],
    ],
    notes="All deployed with security-first defaults: no public endpoints, managed identity, RBAC.",
    col_widths=[2.5, 9.4]
)

# --- SLIDE 9: Architecture - Without Platform LZ ---
make_content_slide(
    "Today's Architecture \u2014 Without Platform LZ",
    [
        "Why standalone for this lab:",
        "\u2022 Faster deployment, self-contained, no hub dependency",
        "\u2022 Same security controls apply \u2014 private endpoints, managed identity",
        "\u2022 In production: enterprise customers use the 'with platform LZ' variant",
        "",
        "Key insight: The architecture is the same.",
        "The difference is whether networking and identity are shared (platform)",
        "or self-contained (standalone).",
    ],
    notes="Link to architecture diagrams in the official repo."
)

# --- SLIDE 10: Deployment Path Comparison ---
make_table_slide(
    "Deployment Path \u2014 Why azd up",
    ["Consideration", "azd up", "Bicep", "Terraform"],
    [
        ["Speed", "~45 min", "~30-60 min", "~30-60 min"],
        ["Setup", "Minimal", "Moderate", "Moderate"],
        ["Customization", "Limited", "Full", "Full"],
        ["Best for", "Labs, PoCs", "Production (Azure-native)", "Production (multi-cloud)"],
    ],
    notes="Today we use azd up for speed. Partners should use Bicep/Terraform for customer deployments.",
    col_widths=[3.0, 2.97, 2.97, 2.97]
)

# --- SLIDE 11: Module 2 divider ---
make_section_slide("2", "Deploy the Landing Zone \U0001f52c HANDS-ON LAB")

# --- SLIDE 12: Lab Overview ---
make_callout_slide(
    "Lab Overview",
    "\u26a0\ufe0f Standard deployment costs ~$2,128-3,098/month. DELETE all resources after the workshop!",
    [
        "Goal: Deploy a complete AI Landing Zone with ~30 Azure resources",
        "Time: ~60 minutes (including deployment wait time)",
        "Steps: Clone \u2192 Configure \u2192 Deploy \u2192 Validate",
        "During deployment (~30-45 min wait): Guided Portal walkthrough",
    ],
    callout_color="orange",
    notes="Emphasize the cost warning. See the official Cost Guide."
)

# --- SLIDE 13: Step 1 - Clone ---
make_code_slide(
    "Step 1 \u2014 Clone the Repository",
    "git clone --recurse-submodules \\\n  https://github.com/microsoft/Deploy-Your-AI-Application-In-Production.git\n\ncd Deploy-Your-AI-Application-In-Production",
    [
        "\u2022 --recurse-submodules is critical \u2014 pulls in Bicep modules",
        "\u2022 Review the repo structure: /infra, /docs, /src",
    ],
    notes="Make sure participants run this exactly as shown."
)

# --- SLIDE 14: Step 2 - Configure ---
make_code_slide(
    "Step 2 \u2014 Initialize and Configure",
    "azd init",
    [
        "Key parameters to set:",
        "\u2022 Environment name: Use workshop naming convention",
        "\u2022 Region: Choose based on quota availability",
        "\u2022 Resource naming prefix: Differentiate from other deployments",
        "\u2022 SKU selections: Balance cost vs. capability",
        "",
        "\U0001f4da Parameter Guide | Roles & Scopes \u2014 see links in chat",
    ],
    notes="Walk through the azd init prompts step by step with participants."
)

# --- SLIDE 15: Step 3 - Deploy ---
make_code_slide(
    "Step 3 \u2014 Deploy",
    "azd up",
    [
        "What's happening behind the scenes (~30-45 min):",
        "\u2022 1. Bicep templates compiled",
        "\u2022 2. Resource group created",
        "\u2022 3. Networking deployed (VNet, subnets, NSGs)",
        "\u2022 4. PaaS services deployed with private endpoints",
        "\u2022 5. RBAC assignments configured",
        "\u2022 6. Application code deployed to Container Apps",
        "",
        "While we wait \u2192 let's explore what's being built...",
    ],
    notes="Start azd up and move to the Portal walkthrough slides while waiting."
)

# --- SLIDE 16: Portal Walkthrough - Networking ---
make_content_slide(
    "Portal Walkthrough \u2014 Networking Layer",
    [
        "Show during deployment wait time:",
        "",
        "\u2022 Virtual network and subnet structure",
        "\u2022 Private endpoints \u2014 no public IP addresses",
        "\u2022 NSG rules \u2014 default deny, explicit allow",
        "\u2022 Private DNS zones for PaaS services",
        "",
        "Design Checklist validation:",
        "\u2705 N-R3: Private Endpoints for all PaaS services",
        "\u2705 N-R4: NSGs on all subnets",
        "\u2705 N-R8: Private DNS zones configured",
    ],
    notes="Walk through these items live in the Azure Portal."
)

# --- SLIDE 17: Portal Walkthrough - AI Services ---
make_content_slide(
    "Portal Walkthrough \u2014 AI Services Layer",
    [
        "Show during deployment wait time:",
        "",
        "\u2022 AI Foundry hub and project structure",
        "\u2022 Azure OpenAI deployments and model endpoints",
        "\u2022 AI Search index configuration",
        "\u2022 Cosmos DB for chat history storage",
        "",
        "Design Checklist validation:",
        "\u2705 D-R1: Thread/chat storage (Cosmos DB)",
        "\u2705 D-R2: File storage (Storage Account)",
    ],
    notes="Show each service in the Portal. Highlight private endpoint configuration."
)

# --- SLIDE 18: Portal Walkthrough - Security ---
make_content_slide(
    "Portal Walkthrough \u2014 Security Layer",
    [
        "Show during deployment wait time:",
        "",
        "\u2022 Managed identities \u2014 no credentials in code",
        "\u2022 Key Vault for secrets management",
        "\u2022 RBAC assignments \u2014 least privilege",
        "\u2022 Defender for Cloud recommendations",
        "",
        "Design Checklist validation:",
        "\u2705 I-R1: Managed Identity throughout",
        "\u2705 I-R3: Entra ID for authentication (not API keys)",
        "\u2705 S-R1: Defender for Cloud enabled",
    ],
    notes="Key point: zero secrets in code. Everything uses Managed Identity."
)

# --- SLIDE 19: Step 4 - Validate ---
make_code_slide(
    "Step 4 \u2014 Validate Deployment",
    "azd show",
    [
        "Verify in the Portal:",
        "\u2610 All resources created in resource group",
        "\u2610 Private endpoints active and connected",
        "\u2610 RBAC assignments correct",
        "\u2610 No public endpoints exposed",
        "\u2610 Diagnostic settings configured",
        "",
        "Troubleshooting: If deployment failed, check quota limits,",
        "region availability, and role assignments. Run 'azd down' and retry.",
    ],
    notes="Give participants 5 minutes to validate their deployment."
)

# --- SLIDE 20: Key Takeaways Module 2 ---
make_content_slide(
    "Key Takeaways \u2014 Module 2",
    [
        "\u2705 A complete production-ready Landing Zone in ~45 minutes",
        "\u2705 Security-first by default (private endpoints, MI, RBAC)",
        "\u2705 Design Checklist items are already implemented in the templates",
        "\u2705 This is what partners should demonstrate to customers",
        "",
        "The infrastructure is deployed. Now let's build a workload on it.",
    ],
    notes="Transition to break."
)

# --- SLIDE 21: Break ---
make_title_slide("\u2615 Break", "15 minutes \u2014 we'll resume with Module 3")

# --- SLIDE 22: Module 3 divider ---
make_section_slide("3", "Configure & Deploy RAG Application \U0001f52c HANDS-ON LAB")

# --- SLIDE 23: RAG Architecture ---
make_content_slide(
    "RAG Application Architecture",
    [
        "User Query \u2192 Container App \u2192 AI Foundry \u2192 AI Search (retrieve)",
        "                                                                  \u2193",
        "                                                    Document Index (embeddings)",
        "                                                                  \u2193",
        "                                             OpenAI (generate) \u2192 Response",
        "",
        "\u2022 Retrieval: AI Search finds relevant document chunks via vector search",
        "\u2022 Augmentation: Retrieved context is added to the LLM prompt",
        "\u2022 Generation: OpenAI model generates the grounded response",
        "",
        "This is deterministic retrieval \u2014 same query returns similar results every time",
    ],
    notes="Draw this flow on the whiteboard if possible."
)

# --- SLIDE 24: Step 1 - Configure AI Foundry ---
make_content_slide(
    "Step 1 \u2014 Configure AI Foundry",
    [
        "Navigate to AI Foundry in the Portal:",
        "",
        "1. Review the project structure (hub \u2192 project)",
        "2. Verify model deployments (GPT-4o, embedding model)",
        "3. Check connections to AI Search, Cosmos DB, Storage",
        "4. Confirm private networking is active (standard setup)",
        "",
        'Key concept: AI Foundry\'s standard setup with private networking = your Landing Zone',
    ],
    notes="Walk through each step live. Have participants follow along."
)

# --- SLIDE 25: Step 2 - Ingest Data ---
make_content_slide(
    "Step 2 \u2014 Ingest Sample Data",
    [
        "1. Upload sample documents to Storage Account",
        "2. Configure AI Search indexer",
        "3. Create embeddings for vector search",
        "4. Verify index population and search results",
        "",
        'Key concept: Quality of retrieval = quality of responses.',
        'Garbage in, garbage out \u2014 even with the best model.',
    ],
    notes="Emphasize data quality. This is where most RAG projects succeed or fail."
)

# --- SLIDE 26: Step 3 - Deploy and Test ---
make_content_slide(
    "Step 3 \u2014 Deploy and Test the Chat App",
    [
        "1. Access the Container App endpoint",
        "2. Test with sample queries",
        "3. Observe the RAG pattern in action:",
        "   \u2022 Query \u2192 Search \u2192 Retrieved context \u2192 Generated response",
        "   \u2022 Responses are grounded in your data, not hallucinated",
        "",
        "Try edge cases:",
        "\u2022 Questions outside the data scope",
        "\u2022 Ambiguous queries",
        "\u2022 Multi-turn conversations (check Cosmos DB chat history)",
    ],
    notes="Give participants 10 minutes to test. Circulate and help."
)

# --- SLIDE 27: Step 4 - Understand What's Happening ---
make_table_slide(
    "Step 4 \u2014 Understand What's Happening",
    ["Metric", "What to Watch"],
    [
        ["Search latency", "Is AI Search returning results fast enough?"],
        ["Token count", "Are prompts too large? Right-size the context window"],
        ["Response quality", "Are retrieved documents relevant?"],
        ["Cost per query", "Prompt + completion tokens \u00d7 model pricing"],
    ],
    notes="Open AI Foundry \u2192 Tracing to see the full request flow. Show token usage, latency breakdown.",
    col_widths=[3.0, 8.9]
)

# --- SLIDE 28: When RAG Is Enough ---
make_callout_slide(
    "When RAG Is Enough",
    '"Not every AI workload needs an agent. Start with RAG and graduate when the use case demands it."',
    [
        "RAG is sufficient when:",
        "\u2705 Customer needs document Q&A, knowledge base, FAQ",
        "\u2705 Inputs are well-defined, outputs are predictable",
        "\u2705 Low latency and cost are priorities",
        "\u2705 No need for tool orchestration or multi-step reasoning",
        "",
        "RAG is the baseline workload \u2014 the one you should lead with",
        "for most customer conversations.",
    ],
    callout_color="blue",
    notes="This is the most important partner talking point from this module."
)

# --- SLIDE 29: Key Takeaways Module 3 ---
make_content_slide(
    "Key Takeaways \u2014 Module 3",
    [
        "\u2705 RAG is the foundation workload on AI Landing Zones",
        "\u2705 Quality of retrieval directly impacts response quality",
        "\u2705 The Landing Zone provides all the infrastructure RAG needs",
        "\u2705 RAG is predictable, cost-effective, and well-understood",
        "\u2705 Know when RAG is enough before recommending agents",
    ],
    notes="Transition to Module 4: From RAG to Agents."
)

# --- SLIDE 30: Module 4 divider ---
make_section_slide("4", "From RAG to Agents")

# --- SLIDE 31: RAG -> Agent Spectrum ---
make_table_slide(
    "The RAG \u2192 Agent Spectrum",
    ["", "Classic RAG", "AI Agent"],
    [
        ["Pipeline", "Fixed: query \u2192 search \u2192 generate", "Dynamic: model decides what to do"],
        ["Tools", "Only retrieval", "Retrieval + actions + APIs + functions"],
        ["Reasoning", "None (template-driven)", "Multi-step, conditional, adaptive"],
        ["Predictability", "High", "Lower (nondeterministic)"],
        ["Cost", "Lower", "Higher (more model calls, tool invocations)"],
        ["Governance", "Simpler", "Complex (tool safety, escalation, audit)"],
    ],
    notes="This table is the foundation for the agent conversation with customers.",
    col_widths=[2.5, 4.7, 4.7]
)

# --- SLIDE 32: When to Graduate to Agents ---
make_two_column_slide(
    "When to Graduate to Agents",
    "Don't Use Agents When",
    [
        "\u2022 Task is structured, predictable, rule-based",
        "  \u2192 Use deterministic code",
        "\u2022 Goal is static knowledge retrieval / Q&A",
        "  \u2192 Use classic RAG",
    ],
    "Use Agents When Task Requires",
    [
        "\u2022 \U0001f504 Dynamic decision-making",
        "  \u2014 multi-step reasoning with conditional logic",
        "\u2022 \U0001f517 Complex orchestration",
        "  \u2014 chaining tools, APIs, services together",
        "\u2022 \U0001f3af Adaptive behavior",
        "  \u2014 ambiguous inputs, intent interpretation",
    ],
    notes="Use the CAF AI Agent Decision Tree to guide this conversation with customers."
)

# --- SLIDE 33: Three Agent Types ---
make_table_slide(
    "Three Agent Types on Your Landing Zone",
    ["Agent Type", "What It Does", "Adds to Your LZ", "Example"],
    [
        ["Productivity", "Retrieve & synthesize", "Minimal \u2014 similar to RAG", "Internal knowledge assistant"],
        ["Action", "Perform specific tasks", "Logic Apps, Functions, APIM", "Ticket creation, data updates"],
        ["Automation", "Multi-step processes", "Orchestration, triggers, governance", "Supply chain optimization"],
    ],
    notes="Same Landing Zone foundation \u2014 different workload complexity on top.",
    col_widths=[2.0, 3.0, 3.5, 3.4]
)

# --- SLIDE 34: Foundry Agent Platform ---
make_table_slide(
    "Microsoft Foundry \u2014 Your Agent Platform",
    ["Agent Build Option", "Description"],
    [
        ["Declarative agents", "Prompt-based, behavior-driven \u2014 simpler to update and version"],
        ["Hosted agents", "Code-first, custom libraries \u2014 full control over logic"],
        ["Workflows", "Orchestrate multi-agent or multi-step processes"],
    ],
    notes="Foundry is what runs on the Landing Zone you just deployed. Standard setup with private networking = your Landing Zone.",
    col_widths=[3.5, 8.4]
)

# --- SLIDE 35: Foundry = Your Landing Zone ---
make_table_slide(
    "Foundry Standard Setup = Your Landing Zone",
    ["Foundry Standard Requirement", "AI Landing Zone Provides"],
    [
        ["Private networking", "\u2705 VNet, Private Endpoints, NSGs"],
        ["Managed Identity", "\u2705 No credentials in code"],
        ["Enterprise data controls", "\u2705 Purview, Key Vault, RBAC"],
        ["AI Search integration", "\u2705 Already deployed"],
        ["Cosmos DB for memory", "\u2705 Already deployed"],
        ["Diagnostic settings", "\u2705 Log Analytics + App Insights"],
    ],
    notes="Key insight: Partners who deploy the Landing Zone are already set up for Foundry agents.",
    col_widths=[4.5, 7.4]
)

# --- SLIDE 36: Single vs Multi-Agent ---
make_table_slide(
    "Single Agent vs. Multi-Agent",
    ["Decision Factor", "Single Agent", "Multi-Agent"],
    [
        ["Complexity", "Lower", "Higher"],
        ["Latency", "Lower", "Higher (handoff overhead)"],
        ["Credential mgmt", "Simpler", "Per-agent credentials"],
        ["State sync", "Built-in", "Must coordinate"],
        ["Cost", "Lower", "Higher"],
    ],
    notes="Rule of thumb: Prototype single, test, then decide. Start with single agent unless you must separate.",
    col_widths=[3.0, 4.45, 4.45]
)

# --- SLIDE 37: Partner Conversation Framework ---
make_content_slide(
    "The Partner Conversation Framework",
    [
        '1. Discovery \u2014 "What problem are you solving? Q&A, task automation, or complex orchestration?"',
        "2. Qualify \u2014 Use the decision tree: deterministic code, RAG, or agent?",
        "3. Scope \u2014 Which agent type? Productivity, Action, or Automation?",
        "4. Architecture \u2014 Single or multi-agent? SaaS or custom build?",
        "5. Infrastructure \u2014 Deploy AI Landing Zone \u2192 configure Foundry standard setup",
        "6. Iterate \u2014 Start simple (RAG or single agent), validate, then graduate",
    ],
    notes="This is the framework partners should use in every customer engagement."
)

# --- SLIDE 38: Key Takeaways Module 4 ---
make_content_slide(
    "Key Takeaways \u2014 Module 4",
    [
        "\u2705 RAG is a subset of what agents can do \u2014 start here",
        "\u2705 Graduate to agents when reasoning and tool orchestration are genuinely needed",
        "\u2705 The Landing Zone you deployed is already agent-ready (Foundry standard setup)",
        "\u2705 Use the decision tree & agent types to qualify customer workloads",
        "\u2705 Start single-agent, prove value, then consider multi-agent",
    ],
    notes="Transition to Module 5: Monitoring."
)

# --- SLIDE 39: Module 5 divider ---
make_section_slide("5", "Monitoring & Observability")

# --- SLIDE 40: Monitoring AI Workloads ---
make_table_slide(
    "Monitoring AI Workloads \u2014 What's Different",
    ["AI-Specific Monitoring", "Why It Matters"],
    [
        ["Token usage & cost tracking", "Predict and control spend"],
        ["Response quality & relevance", "Detect degradation early"],
        ["Hallucination detection", "Trust and safety"],
        ["Model drift over time", "Performance regression"],
        ["RAG retrieval quality", "Right documents found?"],
        ["For agents: tool execution", "Are tools being called correctly?"],
    ],
    notes="Traditional monitoring still applies \u2014 but AI adds new dimensions. Checklist items M-R1 through M-R6.",
    col_widths=[5.0, 6.9]
)

# --- SLIDE 41: What's Already Deployed ---
make_table_slide(
    "What's Already Deployed (from azd up)",
    ["Tool", "What It Monitors"],
    [
        ["Application Insights", "Request tracing, error rates, latency"],
        ["Log Analytics", "Centralized logs, diagnostic settings"],
        ["AI Foundry Tracing", "Model call traces, token usage"],
    ],
    notes="Key metrics: response latency (P50/P95/P99), token consumption, error rate, AI Search performance, Cosmos DB RU consumption.",
    col_widths=[4.0, 7.9]
)

# --- SLIDE 42: AI Foundry Tracing ---
make_content_slide(
    "AI Foundry Tracing Walkthrough",
    [
        "1. Open AI Foundry \u2192 Tracing",
        "2. Find a trace from the RAG query you ran in Module 3",
        "3. Inspect the full chain:",
        "   Search query \u2192 document retrieval \u2192 prompt construction \u2192 model call \u2192 response",
        "4. Identify optimization opportunities:",
        "   \u2022 Slow search? Tune index configuration",
        "   \u2022 Large prompts? Right-size the context window",
        "   \u2022 High token usage? Optimize chunking strategy",
        "",
        "Workshop 3 will go deeper: alerting, dashboards, GenAIOps pipelines",
    ],
    notes="Walk through a live trace if the deployment is complete."
)

# --- SLIDE 43: Key Takeaways Module 5 ---
make_content_slide(
    "Key Takeaways \u2014 Module 5",
    [
        "\u2705 AI workloads need both traditional and AI-specific monitoring",
        "\u2705 azd up deploys monitoring infrastructure by default",
        "\u2705 AI Foundry tracing gives visibility into the full request chain",
        "\u2705 Token usage = cost \u2014 monitor it closely",
        "\u2705 Workshop 3 will cover production monitoring in depth",
    ],
    notes="Transition to Module 6: Wrap-up."
)

# --- SLIDE 44: Module 6 divider ---
make_section_slide("6", "Wrap-up & Next Steps")

# --- SLIDE 45: Workshop Summary ---
make_content_slide(
    "Workshop Summary",
    [
        "Today we accomplished:",
        "",
        "\u2705 Deployed a complete AI Landing Zone (~30 resources)",
        "\u2705 Validated security controls against the Design Checklist",
        "\u2705 Built and tested a RAG chat application",
        "\u2705 Explored when and how to graduate from RAG to AI agents",
        "\u2705 Reviewed monitoring and observability basics",
    ],
    notes="Quick recap before clean-up."
)

# --- SLIDE 46: Clean-up ---
make_callout_slide(
    "\u26a0\ufe0f Clean-up \u2014 DELETE YOUR RESOURCES",
    "Standard deployment costs ~$2,128-3,098/month. Delete everything now unless you are actively experimenting.",
    [
        "azd down",
        "",
        "Set budget alerts if keeping resources:",
        "https://learn.microsoft.com/azure/cost-management-billing/costs/tutorial-acm-create-budgets",
    ],
    callout_color="orange",
    notes="Make sure every participant runs azd down before leaving."
)

# --- SLIDE 47: Key Resources ---
make_table_slide(
    "Key Resources",
    ["Resource", "Link"],
    [
        ["AI Landing Zones Repo", "github.com/Azure/AI-Landing-Zones"],
        ["Deploy-Your-AI-App", "github.com/microsoft/Deploy-Your-AI-App-In-Production"],
        ["Design Checklist", "AI-Landing-Zones-Design-Checklist.md"],
        ["Cost Guide", "AI-Landing-Zones-Cost-Guide.md"],
        ["CAF AI Agent Adoption", "learn.microsoft.com/azure/cloud-adoption-framework/ai-agents/"],
        ["Foundry Agent Quickstart", "learn.microsoft.com/azure/ai-foundry/agents/quickstart"],
    ],
    col_widths=[4.0, 7.9]
)

# --- SLIDE 48: Next Steps ---
make_content_slide(
    "Next Steps",
    [
        "1. Repeat \u2014 Deploy again in your own subscription with different parameters",
        "2. Experiment \u2014 Create a Foundry agent on the deployed infrastructure",
        "3. Read \u2014 Review the full CAF AI Agent Adoption guidance",
        "4. Apply \u2014 Use the deployment + decision framework with your next customer",
        "5. Continue \u2014 Attend Workshop 3: Production Readiness for GenAIOps, CI/CD, scaling",
    ],
    notes="Encourage participants to deploy again on their own."
)

# --- SLIDE 49: Workshop 3 Preview ---
make_content_slide(
    "Workshop 3 Preview \u2014 Production Readiness (GenAIOps Bridge)",
    [
        "\u2022 CI/CD pipelines for AI applications",
        "\u2022 Lifecycle management and model versioning",
        "\u2022 Production monitoring, alerting, and dashboards",
        "\u2022 Cost optimization and scaling strategies",
        "\u2022 Operational runbooks and incident response",
    ],
    notes="Tease Workshop 3 to maintain engagement."
)

# --- SLIDE 50: Q&A ---
make_title_slide(
    "Q&A",
    "Questions? Feedback?",
    "AI Center of Excellence V2 | Partner Enablement Team"
)

# --- SLIDE 51: Thank You ---
make_title_slide(
    "Thank You",
    "Workshop materials available in the partner repository\nFeedback: Survey link in chat\nContact: Arturo Quiroga (PSA) | Anahita Afshari (PSA)",
    "AI Center of Excellence V2 | Q3 FY2026"
)


# ============================================================
# SAVE
# ============================================================

output_path = os.path.join(SCRIPT_DIR, "SLIDE-DECK-EDITABLE.pptx")
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to {output_path}")
