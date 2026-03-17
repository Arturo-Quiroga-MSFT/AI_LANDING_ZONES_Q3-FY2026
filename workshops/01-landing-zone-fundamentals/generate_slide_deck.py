"""
Generate an editable PPTX slide deck for Workshop 1: AI Landing Zones Fundamentals.
Uses python-pptx to create real text, tables, and images (not screenshots).
Output: SLIDE-DECK-EDITABLE.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Microsoft brand colors
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x5A, 0x9E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARNING_ORANGE = RGBColor(0xD8, 0x3B, 0x01)
WARNING_BG = RGBColor(0xFF, 0xF4, 0xCE)
GREEN = RGBColor(0x10, 0x7C, 0x10)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMAGES_DIR = os.path.join(SCRIPT_DIR, "images")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_footer(slide, text="AI Center of Excellence V2 | Q3 FY2026"):
    """Add a small footer to the bottom of a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY


def add_blue_bar(slide, top=Inches(1.35), width=Inches(2.5)):
    """Add a blue accent bar under heading area."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), top, width, Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZURE_BLUE
    bar.line.fill.background()


def set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def make_title_slide(title, subtitle="", footer_text=""):
    """Create a centered title slide with blue accent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, AZURE_BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    if footer_text:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
        tf2 = txBox2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = footer_text
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
        p3.alignment = PP_ALIGN.CENTER

    return slide


def make_section_slide(module_num, title):
    """Create a module divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    # Module number
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Module {module_num}"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xCC, 0xE4, 0xF7)
    p.alignment = PP_ALIGN.CENTER

    # Title
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    return slide


def make_content_slide(title, bullets=None, notes=""):
    """Create a standard content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    # Bullets
    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()

            # Handle bold prefix (text before colon)
            if " \u2014 " in bullet:
                parts = bullet.split(" \u2014 ", 1)
                run1 = p2.add_run()
                run1.text = parts[0] + " \u2014 "
                run1.font.bold = True
                run1.font.size = Pt(18)
                run1.font.color.rgb = DARK_GRAY
                run2 = p2.add_run()
                run2.text = parts[1]
                run2.font.size = Pt(18)
                run2.font.color.rgb = DARK_GRAY
            else:
                p2.text = bullet
                p2.font.size = Pt(18)
                p2.font.color.rgb = DARK_GRAY

            p2.space_before = Pt(8)
            p2.level = 0

    add_footer(slide)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def make_two_column_slide(title, left_title, left_bullets, right_title, right_bullets, notes=""):
    """Create a two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(22)
    lp.font.color.rgb = DARK_BLUE
    lp.font.bold = True
    for b in left_bullets:
        lp2 = ltf.add_paragraph()
        lp2.text = b
        lp2.font.size = Pt(16)
        lp2.font.color.rgb = DARK_GRAY
        lp2.space_before = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(22)
    rp.font.color.rgb = DARK_BLUE
    rp.font.bold = True
    for b in right_bullets:
        rp2 = rtf.add_paragraph()
        rp2.text = b
        rp2.font.size = Pt(16)
        rp2.font.color.rgb = DARK_GRAY
        rp2.space_before = Pt(6)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_table_slide(title, headers, rows, notes="", col_widths=None):
    """Create a slide with a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(11.9)
    table_height = Inches(min(num_rows * 0.5, 5.0))

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.7), Inches(1.6),
        table_width, table_height
    )
    tbl = tbl_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZURE_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = DARK_GRAY

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_image_slide(title, image_path, bullets=None, image_side="right", notes=""):
    """Create a slide with image on one side and text on the other."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    if image_side == "right":
        text_left, text_width = Inches(0.7), Inches(5.5)
        img_left, img_width = Inches(6.5), Inches(6.5)
    else:
        text_left, text_width = Inches(7.0), Inches(5.5)
        img_left, img_width = Inches(0.7), Inches(6.0)

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path, img_left, Inches(1.6), img_width
        )

    # Bullets
    if bullets:
        txBox2 = slide.shapes.add_textbox(text_left, Inches(1.6), text_width, Inches(5.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(16)
            p2.font.color.rgb = DARK_GRAY
            p2.space_before = Pt(8)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_callout_slide(title, callout_text, bullets=None, callout_color="blue", notes=""):
    """Create a slide with a prominent callout box."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = AZURE_BLUE
    p.font.bold = True

    add_blue_bar(slide)

    # Callout box
    bar_color = AZURE_BLUE if callout_color == "blue" else WARNING_ORANGE
    bg_color = RGBColor(0xE8, 0xF4, 0xFD) if callout_color == "blue" else WARNING_BG

    callout_top = Inches(1.6)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), callout_top, Inches(11.9), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), callout_top, Inches(0.06), Inches(1.2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = bar_color
    accent.line.fill.background()

    ctxBox = slide.shapes.add_textbox(Inches(1.0), callout_top + Inches(0.15), Inches(11.3), Inches(0.9))
    ctf = ctxBox.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = callout_text
    cp.font.size = Pt(16)
    cp.font.color.rgb = DARK_GRAY
    cp.font.italic = True

    # Bullets below
    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(11.9), Inches(4.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(18)
            p2.font.color.rgb = DARK_GRAY
            p2.space_before = Pt(8)

    add_footer(slide)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# ============================================================
# SLIDE GENERATION
# ============================================================

# --- SLIDE 1: Title ---
make_title_slide(
    "AI Landing Zones Fundamentals",
    "Partner Enablement Workshop",
    "AI Center of Excellence V2 | Q3 FY2026"
)

# --- SLIDE 2: Agenda ---
make_table_slide(
    "Agenda",
    ["Module", "Topic", "Duration"],
    [
        ["1", "What is AI Landing Zone?", "25 min"],
        ["2", "Reference Architectures", "30 min"],
        ["\u2014", "Break", "10 min"],
        ["3", "Design Checklist Walkthrough", "40 min"],
        ["4", "Deployment Options", "25 min"],
        ["5", "Partner Engagement Scenarios", "20 min"],
        ["\u2014", "Wrap-up & Next Steps", "15 min"],
    ],
    notes="Walk through the agenda. Each module builds on the previous one.",
    col_widths=[1.5, 7.0, 3.4]
)

# --- SLIDE 3: Learning Objectives ---
make_content_slide(
    "Learning Objectives",
    [
        "Explain what Azure AI Landing Zones are and why they matter",
        "Identify key architecture components",
        "Compare deployment options (with/without Platform LZ)",
        "Navigate the official Design Checklist across 10 design areas",
        "Choose appropriate deployment approach for different scenarios",
        "Plan customer engagements using partner frameworks",
    ],
    notes="These objectives are measurable. We'll check for understanding at the end of each module."
)

# --- SLIDE 4: Housekeeping ---
make_callout_slide(
    "Housekeeping",
    "\u26a0\ufe0f Cost Warning: If you deploy Azure resources, a full AI Landing Zone can cost $1,500-3,500+/month. Delete all resources immediately after. See Prerequisites doc for details.",
    [
        "Duration: 2-3 hours",
        "Breaks: 10 min after Module 2",
        "Questions: Welcome anytime \u2014 raise hand or use chat",
        "Materials: Links shared in chat",
        "Recording: [Yes/No]",
    ],
    callout_color="orange",
    notes="Emphasize the cost warning early. Participants who plan to deploy resources should set up budget alerts first."
)

# --- SLIDE 5: Module 1 divider ---
make_section_slide("1", "What is AI Landing Zone?")

# --- SLIDE 6: The AI Challenge ---
make_callout_slide(
    "The AI Challenge",
    "How do we accelerate production-ready AI deployments?",
    [
        "Enterprises want to deploy Gen AI workloads",
        "Production requires: security, governance, scalability",
        "Starting from scratch is slow and risky",
        "Partners need a repeatable, validated approach",
    ],
    callout_color="blue",
    notes="Ask: 'How many have had customers stall on AI deployment when security comes up?'"
)

# --- SLIDE 7: AI Landing Zone Definition ---
make_callout_slide(
    "AI Landing Zone \u2014 Defined",
    '"A secure, resilient, and scalable reference architecture for AI apps and agents"',
    [
        "Purpose-built for Azure AI Foundry and Gen AI workloads",
        "Implements Azure best practices from Day 0",
        "Available as: Bicep, Terraform, Portal (coming soon)",
        "Working IaC code \u2014 not just documentation",
        "Source: github.com/Azure/AI-Landing-Zones",
    ],
    notes="This is hosted on the Azure GitHub org \u2014 official Microsoft guidance, not a community project."
)

# --- SLIDE 8: Where AI LZ Fits in CAF ---
make_image_slide(
    "Where AI Landing Zones Fit in CAF",
    os.path.join(IMAGES_DIR, "azure-landing-zone-conceptual-diagram.png"),
    [
        "Platform LZ \u2192 Shared services (Identity, Connectivity, Management)",
        "Application LZ \u2192 Where workloads deploy",
        "AI Landing Zone \u2192 Goes here \u261d\ufe0f",
        "",
        "AI is just another workload from CAF perspective.",
        "NOT a separate landing zone type.",
    ],
    notes="Critical concept: AI Landing Zones are Application Landing Zones within CAF."
)

# --- SLIDE 9: Platform vs Application LZ ---
make_two_column_slide(
    "Platform vs. Application Landing Zone",
    "Platform Team Provides",
    [
        "\u2022 Hub virtual network",
        "\u2022 Azure Firewall, VPN/ExpressRoute",
        "\u2022 Azure Bastion for jump box access",
        "\u2022 Private DNS zones",
        "\u2022 DDoS protection",
        "\u2022 Governance policies at MG level",
    ],
    "Workload Team Owns (AI LZ)",
    [
        "\u2022 Spoke virtual network subnets, NSGs",
        "\u2022 AI Foundry, AI Search, Cosmos DB",
        "\u2022 App Service / Container Apps",
        "\u2022 Private endpoints to PaaS services",
        "\u2022 Workload-specific monitoring",
    ],
    notes="Partner Question: Does your customer have an existing Platform LZ?"
)

# --- SLIDE 10: CAF & WAF Alignment ---
make_two_column_slide(
    "CAF & WAF Alignment",
    "Cloud Adoption Framework",
    [
        "\u2022 AI Landing Zone = 'AI Ready' stage",
        "\u2022 Specifically: 'AI on Azure Platforms (PaaS)'",
        "\u2022 Links to CAF AI Strategy and Checklist",
    ],
    "Well-Architected Framework",
    [
        "\u2022 Reliability \u2014 Resilient deployments",
        "\u2022 Security \u2014 Zero-trust by design",
        "\u2022 Cost \u2014 Optimization built in",
        "\u2022 Operations \u2014 Observable by default",
        "\u2022 Performance \u2014 Scalable architecture",
    ],
    notes="This integrates with the broader cloud adoption story customers already know."
)

# --- SLIDE 11: Use Cases & Agent Types ---
make_table_slide(
    "Supported Use Cases & AI Agent Types",
    ["Agent Type", "What It Does", "Example"],
    [
        ["Productivity", "Retrieve & synthesize information", "Knowledge management, customer support"],
        ["Action", "Perform specific tasks in workflows", "Ticket creation, system monitoring"],
        ["Automation", "Multi-step processes, minimal oversight", "Supply chain optimization, approval flows"],
    ],
    notes="Classic RAG = deterministic retrieval (chat, Q&A). AI Agents = adaptive reasoning. Key question: does your customer need classic RAG or AI agents?",
    col_widths=[2.5, 4.5, 4.9]
)

# --- SLIDE 12: Agent Architecture ---
make_content_slide(
    "AI Agent Architecture \u2014 5 Core Components",
    [
        "1. Generative AI Model \u2014 The reasoning engine",
        "2. Instructions \u2014 Scope, boundaries, behavioral guidelines",
        "3. Retrieval (Knowledge) \u2014 Grounding data and context (reduces hallucinations)",
        "4. Actions (Tools) \u2014 Functions, APIs, systems the agent can call",
        "5. Memory \u2014 Conversation history and state",
        "",
        "Key Insight: Components 3 (Retrieval) and 4 (Actions) are what the",
        "Landing Zone infrastructure directly supports \u2014 AI Search, Cosmos DB,",
        "APIs, Logic Apps, Functions.",
    ],
    notes="When we look at the reference architecture, you'll see how each LZ component maps to these agent building blocks."
)

# --- SLIDE 13: Agents vs Classic RAG ---
make_two_column_slide(
    "When to Use Agents vs. Classic RAG",
    "\u274c Don\u2019t Use Agents When",
    [
        "\u2022 Task is structured, predictable, rule-based \u2192 Use deterministic code",
        "\u2022 Goal is static knowledge retrieval / Q&A \u2192 Use classic RAG",
    ],
    "\u2705 Use Agents When",
    [
        "\u2022 Dynamic decision-making (multi-step reasoning)",
        "\u2022 Complex orchestration (chaining tools, APIs)",
        "\u2022 Adaptive behavior (ambiguous inputs, intent interpretation)",
    ],
    notes="Partner talking point: 'Not every AI workload needs an agent. Start simple. Classic RAG is cheaper, faster, and more predictable.'"
)

# --- SLIDE 14: Key Takeaways Module 1 ---
make_content_slide(
    "Key Takeaways \u2014 Module 1",
    [
        "AI Landing Zone = production-ready AI foundation",
        "Supports both classic RAG and AI agent workloads",
        "Aligns with CAF (including AI Agent Adoption) and WAF",
        "Accelerates enterprise AI adoption",
        "Reduces risk and time-to-production",
        "Use the agent decision tree to help partners qualify workload type",
    ],
    notes="Quick check: 'Can someone explain what makes AI Landing Zone different from just deploying Azure AI services directly?'"
)

# --- SLIDE 15: Module 2 divider ---
make_section_slide("2", "Reference Architectures")

# --- SLIDE 16: Two Architecture Options ---
make_table_slide(
    "Two Architecture Options",
    ["", "With Platform LZ", "Without Platform LZ"],
    [
        ["Target", "Enterprise", "Greenfield / PoC"],
        ["Networking", "Hub-spoke (shared)", "Self-contained"],
        ["Security", "Central Firewall, Bastion", "Own networking & security"],
        ["Speed", "More setup", "Faster start"],
        ["Production-ready", "\u2705 Yes", "\u2705 Yes"],
    ],
    notes="Both are valid. Choose based on customer context. The key question: 'Does the customer have an existing Azure Landing Zone?'",
    col_widths=[3.0, 4.45, 4.45]
)

# --- SLIDE 17: With Platform LZ ---
make_image_slide(
    "With Platform Landing Zone",
    os.path.join(IMAGES_DIR, "AI-Landing-Zone-with-platform.png"),
    [
        "Leverages existing hub-spoke networking",
        "Uses central Firewall, Bastion, DNS",
        "AI workload is an application landing zone subscription",
        "Best for: enterprises, regulated industries",
    ],
    notes="Walk through layers top to bottom. Notice all private endpoints \u2014 no public exposure."
)

# --- SLIDE 18: Without Platform LZ ---
make_image_slide(
    "Without Platform Landing Zone",
    os.path.join(IMAGES_DIR, "AI-Landing-Zone-without-platform.png"),
    [
        "Self-contained in a single subscription",
        "Includes own networking, security, governance",
        "Faster to deploy for isolated scenarios",
        "Best for: PoCs, smaller orgs, quick starts",
        "Can migrate to platform model later",
    ],
    notes="A 'complete' landing zone in miniature. Not less secure \u2014 just manages everything internally."
)

# --- SLIDES 19-23: Component slides ---
for comp_title, comp_rows in [
    ("Key Components \u2014 AI Layer", [
        ["Azure AI Foundry", "Development platform \u2014 model management, prompt engineering"],
        ["AI Services", "Model endpoints (GPT-4, embeddings, etc.)"],
        ["AI Search", "RAG retrieval \u2014 vector search, semantic ranking"],
        ["Connections", "Service integrations and data source links"],
    ]),
    ("Key Components \u2014 Compute Layer", [
        ["Container Apps", "Microservices runtime \u2014 API backends"],
        ["App Service", "Web front-ends \u2014 user-facing apps"],
        ["Functions", "Event-driven compute \u2014 triggers, processing"],
        ["Scaling", "Auto-scale capabilities across all compute"],
    ]),
    ("Key Components \u2014 Data Layer", [
        ["Cosmos DB", "Chat history, context storage, agent memory"],
        ["Storage Account", "Document storage, files, blobs"],
        ["AI Search", "Vector search, embeddings index"],
        ["(Optional) Fabric", "Enterprise data foundation, analytics"],
    ]),
    ("Key Components \u2014 Security Layer", [
        ["Private Endpoints", "No public exposure for PaaS services"],
        ["Key Vault", "Secrets and certificate management"],
        ["Managed Identity", "No credentials in code \u2014 zero secrets"],
        ["NSGs & Firewall", "Network isolation and egress control"],
        ["Bastion", "Secure admin access to jump boxes"],
    ]),
    ("Key Components \u2014 Governance Layer", [
        ["Defender for Cloud", "Security posture management"],
        ["Purview", "Data governance and classification"],
        ["Azure Policy", "Compliance automation and enforcement"],
        ["Azure Monitor", "Observability \u2014 logs, metrics, alerts"],
    ]),
]:
    make_table_slide(comp_title, ["Component", "Purpose"], comp_rows, col_widths=[3.5, 8.4])

# --- SLIDE 24: Choosing Architecture ---
make_table_slide(
    "Choosing Your Architecture",
    ["Factor", "With Platform LZ", "Without Platform LZ"],
    [
        ["Customer has existing ALZ", "\u2705 Recommended", "\u26a0\ufe0f Consider migration"],
        ["Greenfield deployment", "\u26a0\ufe0f More setup", "\u2705 Faster start"],
        ["Enterprise / regulated", "\u2705 Best fit", "\u26a0\ufe0f May need upgrade"],
        ["PoC / Pilot", "\u26a0\ufe0f Overkill", "\u2705 Right-sized"],
    ],
    col_widths=[4.0, 3.95, 3.95]
)

# --- SLIDE 25: Key Takeaways Module 2 ---
make_content_slide(
    "Key Takeaways \u2014 Module 2",
    [
        "Two validated architecture options",
        "Choose based on customer context",
        "All components are Azure PaaS",
        "Security-first design by default",
        "Both options are production-ready",
    ],
    notes="Break time. When we come back: the Design Checklist."
)

# --- SLIDE 26: Break ---
make_title_slide("\u2615 Break", "10 minutes")

# --- SLIDE 27: Module 3 divider ---
make_section_slide("3", "Design Checklist Walkthrough")

# --- SLIDE 28: The Design Checklist ---
make_callout_slide(
    "The Design Checklist",
    "Your new best friend for customer engagements.",
    [
        "40+ recommendations across 10 design areas",
        "Based on CAF + WAF best practices",
        "Works for greenfield and brownfield",
        "Assessment tool, design guide, and conversation framework",
        "",
        "Link: github.com/Azure/AI-Landing-Zones \u2192 docs/AI-Landing-Zones-Design-Checklist.md",
    ],
    notes="Open the checklist in a browser tab. Show the actual document."
)

# --- SLIDE 29: 10 Design Areas ---
make_table_slide(
    "10 Design Areas Overview",
    ["Area", "# Recs", "Key Focus"],
    [
        ["Networking", "9", "Private endpoints, Firewall, DNS"],
        ["Identity", "6", "Managed Identity, Entra ID, RBAC"],
        ["Security", "5", "Defender, Purview, Content Safety"],
        ["Monitoring", "6", "Tracing, alerts, drift detection"],
        ["Cost", "4", "Pricing models, auto-shutdown"],
        ["Data", "3", "Storage, data governance"],
        ["Governance", "5", "Policy, compliance, Responsible AI"],
        ["Reliability", "1", "Resilience patterns"],
        ["Resource Org", "4", "Subscription and resource structure"],
        ["Compute", "1", "Compute selection"],
    ],
    notes="We'll deep-dive into Networking, Identity, Security, and Monitoring first.",
    col_widths=[2.5, 1.5, 7.9]
)

# --- SLIDES 30-35: Design area deep dives ---
for area_title, recs in [
    ("Deep Dive \u2014 Networking (N-R1 to N-R9)", [
        ["N-R1", "DDoS Protection", "High"],
        ["N-R2", "Bastion for jumpbox access", "High"],
        ["N-R3", "Private endpoints everywhere", "Critical"],
        ["N-R4", "NSGs on all subnets", "High"],
        ["N-R5", "WAF via App Gateway/Front Door", "High"],
        ["N-R6", "APIM as AI Gateway", "Critical"],
        ["N-R7", "Firewall for egress control", "High"],
        ["N-R8", "Private DNS zones", "High"],
        ["N-R9", "Restrict outbound by default", "Critical"],
    ]),
    ("Deep Dive \u2014 Identity (I-R1 to I-R6)", [
        ["I-R1", "Managed Identity with least privilege", "Critical"],
        ["I-R2", "MFA + PIM for sensitive accounts", "High"],
        ["I-R3", "Entra ID for authentication (not API keys)", "Critical"],
        ["I-R4", "Conditional Access policies", "High"],
        ["I-R5", "RBAC with least privilege", "High"],
        ["I-R6", "Disable key-based access entirely", "Critical"],
    ]),
    ("Deep Dive \u2014 Security (S-R1 to S-R5)", [
        ["S-R1", "Defender for Cloud recommendations", "High"],
        ["S-R2", "Microsoft Cloud Security Baseline", "High"],
        ["S-R3", "Purview for data protection", "High"],
        ["S-R4", "MITRE ATLAS + OWASP for AI risks", "Critical"],
        ["S-R5", "AI Content Safety for outputs", "Critical"],
    ]),
    ("Deep Dive \u2014 Monitoring (M-R1 to M-R6)", [
        ["M-R1", "Monitor models, resources, data", "High"],
        ["M-R2", "Azure Monitor Baseline Alerts", "High"],
        ["M-R3", "AI Foundry tracing and evaluation", "Critical"],
        ["M-R4", "Diagnostic settings to Log Analytics", "High"],
        ["M-R5", "Model and data drift detection", "Critical"],
        ["M-R6", "Network Watcher troubleshooting", "Medium"],
    ]),
    ("Deep Dive \u2014 Cost (CO-R1 to CO-R4)", [
        ["CO-R1", "Understand AI Foundry pricing", "Critical"],
        ["CO-R2", "Combine PTU + PAYGO endpoints", "Critical"],
        ["CO-R3", "Consider global deployment types", "High"],
        ["CO-R4", "Auto-shutdown for non-prod", "High"],
    ]),
    ("Deep Dive \u2014 Governance (G-R1 to G-R5)", [
        ["G-R1", "Built-in AI-related Azure policies", "High"],
        ["G-R2", "Regulatory compliance (NIST AI RMF)", "High"],
        ["G-R3", "Responsible AI dashboard", "Critical"],
        ["G-R4", "AI Content Safety for testing", "High"],
        ["G-R5", "Policy to govern model catalog", "High"],
    ]),
]:
    make_table_slide(area_title, ["Rec", "Recommendation", "Priority"], recs, col_widths=[1.5, 7.5, 2.9])

# --- SLIDE 36: How to use the checklist ---
make_content_slide(
    "How to Use the Checklist",
    [
        "1. Assessment \u2014 Walk through with customer current state",
        "2. Gap Analysis \u2014 Identify what's missing or non-compliant",
        "3. Prioritization \u2014 P0 (must-have), P1 (should-have), P2 (nice-to-have)",
        "4. Implementation \u2014 Use IaC templates to deploy",
        "5. Validation \u2014 Re-run checklist post-deployment",
        "",
        "This becomes your engagement framework for every customer conversation.",
    ],
    notes="If time allows: have participants pick one area and identify the most commonly missed recommendation."
)

# --- SLIDE 37: Key Takeaways Module 3 ---
make_content_slide(
    "Key Takeaways \u2014 Module 3",
    [
        "Design Checklist is your primary tool",
        "10 areas cover comprehensive requirements",
        "Use with every customer engagement",
        "Reference, don't recreate \u2014 point to official docs",
    ]
)

# --- SLIDE 38: Module 4 divider ---
make_section_slide("4", "Deployment Options")

# --- SLIDE 39: Four Deployment Paths ---
make_table_slide(
    "Four Deployment Paths",
    ["Path", "Tool", "Speed", "Customization"],
    [
        ["A", "azd up (Azure Developer CLI)", "~45 min", "Limited"],
        ["B", "Bicep", "~30-60 min", "Full"],
        ["C", "Terraform", "~30-60 min", "Full"],
        ["D", "Portal", "Manual", "Limited"],
    ],
    col_widths=[1.5, 4.5, 2.95, 2.95]
)

# --- SLIDE 40: Path A - azd up ---
make_content_slide(
    "Path A \u2014 Azure Developer CLI (azd up)",
    [
        "The fastest path to production:",
        "",
        "    git clone --recurse-submodules <repo>",
        "    azd up",
        "",
        "30+ resources deployed in ~45 minutes",
        "AI Foundry, Fabric, Purview, AI Search \u2014 all configured",
        "Best for: PoCs, demos, fast validation",
        "",
        "Repo: github.com/microsoft/Deploy-Your-AI-Application-In-Production",
    ],
    notes="This is the 'show me it works' path. Great for demos and initial customer validation."
)

# --- SLIDE 41: Cost Warning ---
make_table_slide(
    "\u26a0\ufe0f Cost Warning \u2014 Know Before You Deploy",
    ["Service", "Est. Monthly Cost"],
    [
        ["Azure Firewall (Standard)", "$900-1,250"],
        ["Azure DDoS Protection (Standard)", "$2,944"],
        ["Azure Bastion (Standard)", "$140-200"],
        ["Azure AI Search", "$75-300+"],
        ["Azure Cosmos DB", "$25-500+"],
        ["\u2014\u2014\u2014", "\u2014\u2014\u2014"],
        ["Full 'with Platform LZ'", "$1,500-3,500+/month"],
        ["'Without Platform LZ'", "$300-1,200+/month"],
    ],
    notes="CRITICAL: Do NOT let participants leave resources running. Remind them to run 'azd down' or delete the resource group.",
    col_widths=[7.0, 4.9]
)

# --- SLIDE 42: Bicep vs Terraform ---
make_two_column_slide(
    "Path B vs. Path C \u2014 Bicep vs. Terraform",
    "Bicep",
    [
        "\u2022 Azure-native, first-class support",
        "\u2022 No state file management",
        "\u2022 Best Azure feature coverage",
        "\u2022 AVM (Azure Verified Modules) based",
        "\u2022 Choose if: Azure-only, Microsoft-aligned",
    ],
    "Terraform",
    [
        "\u2022 Multi-cloud capable",
        "\u2022 Flexible state management",
        "\u2022 Mature ecosystem",
        "\u2022 AVM (Azure Verified Modules) based",
        "\u2022 Choose if: Multi-cloud, existing TF investment",
    ],
    notes="Both use Azure Verified Modules \u2014 same foundation, different syntax."
)

# --- SLIDE 43: Decision Framework ---
make_content_slide(
    "Decision Framework",
    [
        "Need fastest deployment?",
        "    \u2514\u2500 Yes \u2192 azd up",
        "",
        "Multi-cloud strategy?",
        "    \u2514\u2500 Yes \u2192 Terraform",
        "",
        "Azure-native team?",
        "    \u2514\u2500 Yes \u2192 Bicep",
        "",
        "No IaC expertise?",
        "    \u2514\u2500 azd up (or wait for Portal)",
    ],
    notes="See IaC Decision Framework doc for details."
)

# --- SLIDE 44: Key Takeaways Module 4 ---
make_content_slide(
    "Key Takeaways \u2014 Module 4",
    [
        "Multiple valid deployment paths",
        "Choose based on customer context",
        "All paths are production-ready",
        "\u26a0\ufe0f Watch costs \u2014 delete resources after learning",
        "See IaC Decision Framework for detailed guidance",
    ]
)

# --- SLIDE 45: Module 5 divider ---
make_section_slide("5", "Partner Engagement Scenarios")

# --- SLIDES 46-50: Scenario slides ---
for s_title, s_context, s_rec, s_focus, s_questions in [
    (
        "Scenario 1 \u2014 Enterprise Customer",
        "Large org with existing Azure Landing Zones",
        "Bicep/Terraform + Platform LZ architecture",
        "Networking, Identity integration",
        "Current ALZ config? Governance policies? Integration needs?"
    ),
    (
        "Scenario 2 \u2014 Greenfield / PoC",
        "New to Azure or exploring AI",
        "azd up (Deploy-Your-AI-App)",
        "Fast deployment, validation",
        "Timeline? Expansion plans? Who needs access?"
    ),
    (
        "Scenario 3 \u2014 Regulated Industry",
        "Healthcare, finance, government",
        "Bicep + extensive security review",
        "All Security (S-R1-5), Governance (G-R1-5)",
        "Compliance frameworks? Data residency? Audit schedule?"
    ),
    (
        "Scenario 4 \u2014 Cost-Sensitive",
        "Limited budget, cost optimization critical",
        "Minimal deployment + Cost checklist",
        "Cost (CO-R1-4), Compute (C-R1)",
        "Budget? Workload predictability? Essential vs. nice-to-have?"
    ),
]:
    make_table_slide(
        s_title,
        ["Aspect", "Details"],
        [
            ["Context", s_context],
            ["Recommendation", s_rec],
            ["Focus Areas", s_focus],
            ["Key Questions", s_questions],
        ],
        col_widths=[3.0, 8.9]
    )

# --- SLIDE 51: Scenario 5 - AI Agents ---
make_table_slide(
    "Scenario 5 \u2014 Customer Building AI Agents",
    ["Agent Type", "LZ Infrastructure Needs", "Additional Components"],
    [
        ["Productivity", "Standard (AI Search, Cosmos DB, Foundry)", "Minimal \u2014 similar to RAG"],
        ["Action", "Standard + APIs/integrations", "Logic Apps, Functions, APIM"],
        ["Automation", "Full + orchestration + governance", "Workflows, triggers, audit trails"],
    ],
    notes="Key questions: What tasks will the agent perform? Single or multi-agent? Who owns governance?",
    col_widths=[2.5, 5.0, 4.4]
)

# --- SLIDE 52: Key Takeaways Module 5 ---
make_content_slide(
    "Key Takeaways \u2014 Module 5",
    [
        "Tailor approach to customer context",
        "Use Design Checklist as conversation guide",
        "Different scenarios need different paths",
        "Ask discovery questions upfront",
        "Use the CAF agent decision tree to qualify agent workloads",
    ]
)

# --- SLIDE 53: Wrap-up divider ---
make_title_slide("Wrap-Up", "")

# --- SLIDE 54: Workshop Summary ---
make_content_slide(
    "Workshop Summary",
    [
        "\u2705 What AI Landing Zones are and why they matter",
        "\u2705 Classic RAG vs. AI Agents \u2014 when to use each",
        "\u2705 Two reference architecture options",
        "\u2705 10 design areas and key recommendations",
        "\u2705 Four deployment paths",
        "\u2705 Five partner engagement scenarios (including AI agents)",
    ]
)

# --- SLIDE 55: Key Resources ---
make_table_slide(
    "Key Resources",
    ["Resource", "Link"],
    [
        ["AI Landing Zones Repo", "github.com/Azure/AI-Landing-Zones"],
        ["Deploy-Your-AI-App", "github.com/microsoft/Deploy-Your-AI-Application-In-Production"],
        ["CAF AI Agent Adoption", "learn.microsoft.com/azure/cloud-adoption-framework/ai-agents/"],
        ["Diagram Builder", "github.com/Arturo-Quiroga-MSFT/azure-architecture-diagram-builder"],
        ["Design Checklist", "github.com/Azure/AI-Landing-Zones \u2192 docs/"],
    ],
    col_widths=[4.0, 7.9]
)

# --- SLIDE 56: Next Steps ---
make_content_slide(
    "Next Steps",
    [
        "🧪 Practice \u2014 Deploy a test environment (see cost warnings!)",
        "\u26a0\ufe0f Delete resources \u2014 Tear down immediately after practice",
        "Read \u2014 Review Design Checklist in detail",
        "Explore \u2014 Review CAF AI Agent Adoption guidance",
        "Apply \u2014 Use with your next customer engagement",
        "Continue \u2014 Attend Workshop 2 (hands-on deployment)",
    ]
)

# --- SLIDE 57: Workshop 2 Preview ---
make_content_slide(
    "Workshop 2 Preview \u2014 From RAG to Agents",
    [
        "Hands-on lab: Deploy Landing Zone + RAG chat application",
        "Configure AI Foundry with standard (private) setup",
        "Understand when & how to graduate from RAG to agents",
        "Explore Microsoft Foundry agent capabilities",
        "Implement monitoring and observability",
    ]
)

# --- SLIDE 58: Q&A ---
make_title_slide("Q&A", "Questions? Feedback?")

# --- SLIDE 59: Thank You ---
make_title_slide(
    "Thank You",
    "AI Center of Excellence V2\nPartner Enablement Team",
    "Q3 FY2026"
)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(SCRIPT_DIR, "SLIDE-DECK-EDITABLE.pptx")
prs.save(output_path)
print(f"Editable PPTX saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
