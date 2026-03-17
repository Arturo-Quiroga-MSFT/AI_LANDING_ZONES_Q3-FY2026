#!/usr/bin/env python3
"""
Analyze existing PPTX presentations and map reusable slides to Workshop modules.

Usage:
    python scripts/analyze_pptx_for_workshop.py

Output:
    scripts/output/workshop1-slide-reuse-map.md
"""

import os
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# ──────────────────────────────────────────────────────────────────────
# Workshop 1 modules and their topic keywords for matching
# ──────────────────────────────────────────────────────────────────────
WORKSHOP1_MODULES = {
    "Module 1: What is AI Landing Zone?": {
        "keywords": [
            "landing zone", "what is", "caf", "cloud adoption framework",
            "application landing zone", "platform landing zone", "ai landing",
            "waf", "well-architected", "use case", "agent type",
            "productivity agent", "action agent", "automation agent",
            "ai agent", "rag", "retrieval augmented", "decision tree",
            "when to use agent", "when not to use",
        ],
        "slides": [],  # will be populated
    },
    "Module 2: Reference Architectures": {
        "keywords": [
            "reference architecture", "with platform", "without platform",
            "hub-spoke", "hub spoke", "architecture diagram", "architecture overview",
            "ai foundry", "ai services", "cosmos db", "ai search",
            "private endpoint", "key vault", "managed identity",
            "container app", "app service", "bastion", "firewall",
            "nsg", "dns", "network", "compute layer", "data layer",
            "security layer", "governance layer", "purview", "defender",
            "azure monitor", "components",
        ],
        "slides": [],
    },
    "Module 3: Design Checklist Walkthrough": {
        "keywords": [
            "design checklist", "design area", "recommendation",
            "networking", "identity", "security", "monitoring", "cost",
            "governance", "reliability", "resource org", "compute",
            "ddos", "rbac", "managed identity", "conditional access",
            "defender for cloud", "azure policy", "responsible ai",
            "content safety", "model drift", "data drift",
            "n-r1", "i-r1", "s-r1", "m-r1", "co-r1", "g-r1",
            "ptu", "paygo",
        ],
        "slides": [],
    },
    "Module 4: Deployment Options": {
        "keywords": [
            "deployment", "deploy", "azd", "bicep", "terraform",
            "portal", "deploy to azure", "iac", "infrastructure as code",
            "azure developer cli", "azd up", "azure verified module",
            "avm", "decision framework", "cost", "pricing",
            "cost guide", "budget",
        ],
        "slides": [],
    },
    "Module 5: Partner Engagement Scenarios": {
        "keywords": [
            "partner", "engagement", "scenario", "enterprise",
            "greenfield", "poc", "regulated", "cost-sensitive",
            "next step", "call to action", "workshop 2", "workshop 3",
            "customer conversation", "qualification",
        ],
        "slides": [],
    },
}


def extract_text_from_slide(slide):
    """Extract all text content from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return texts


def get_slide_title(slide):
    """Get the title of a slide."""
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()
    # Fallback: look for the largest text or first text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) < 200:
                return text[:120]
    return "(No title)"


def score_slide_for_module(slide_text_lower, module_keywords):
    """Score how relevant a slide is to a module based on keyword matches."""
    score = 0
    matched = []
    for kw in module_keywords:
        if kw in slide_text_lower:
            score += 1
            matched.append(kw)
    return score, matched


def analyze_presentation(pptx_path):
    """Analyze a single PPTX file and return slide data."""
    short_name = os.path.basename(pptx_path)
    print(f"  Analyzing: {short_name} ...", flush=True)
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ⚠️  Failed to open {short_name}: {e}")
        return []

    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        title = get_slide_title(slide)
        texts = extract_text_from_slide(slide)
        full_text = " ".join(texts)

        # Count non-empty shapes with text
        text_shapes = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        has_table = any(s.has_table for s in slide.shapes)
        try:
            has_image = any(
                s.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
                for s in slide.shapes
            )
        except (NotImplementedError, AttributeError):
            has_image = False

        slides_data.append({
            "file": short_name,
            "slide_num": i,
            "title": title,
            "full_text": full_text,
            "text_preview": full_text[:300],
            "text_shapes": text_shapes,
            "has_table": has_table,
            "has_image": has_image,
            "word_count": len(full_text.split()),
        })

    print(f"    → {len(slides_data)} slides extracted")
    return slides_data


def map_slides_to_modules(all_slides):
    """Map each slide to Workshop 1 modules based on keyword relevance."""
    # Reset
    for mod in WORKSHOP1_MODULES.values():
        mod["slides"] = []

    for slide in all_slides:
        text_lower = slide["full_text"].lower()
        best_module = None
        best_score = 0
        best_matched = []

        for mod_name, mod_data in WORKSHOP1_MODULES.items():
            score, matched = score_slide_for_module(text_lower, mod_data["keywords"])
            if score > best_score:
                best_score = score
                best_module = mod_name
                best_matched = matched

        if best_score >= 2:  # At least 2 keyword matches
            slide["relevance_score"] = best_score
            slide["matched_keywords"] = best_matched
            WORKSHOP1_MODULES[best_module]["slides"].append(slide)


def generate_report(all_slides, output_path):
    """Generate the Workshop 1 slide reuse mapping report."""
    map_slides_to_modules(all_slides)

    lines = []
    lines.append("# Workshop 1: Slide Reuse Map from Existing Presentations\n")
    lines.append(f"**Generated**: {__import__('datetime').date.today()}\n")
    lines.append("**Purpose**: Identify reusable slides from existing PPTX decks for Workshop 1 modules\n")
    lines.append("---\n")

    # Summary table
    lines.append("## Summary\n")
    lines.append("| Source Deck | Total Slides | Relevant to WS1 |")
    lines.append("|------------|-------------|-----------------|")

    file_stats = {}
    for slide in all_slides:
        fname = slide["file"]
        if fname not in file_stats:
            file_stats[fname] = {"total": 0, "relevant": 0}
        file_stats[fname]["total"] += 1

    # Count relevant
    for mod_data in WORKSHOP1_MODULES.values():
        for slide in mod_data["slides"]:
            file_stats[slide["file"]]["relevant"] += 1

    for fname, stats in file_stats.items():
        lines.append(f"| {fname} | {stats['total']} | {stats['relevant']} |")
    lines.append("")

    # Per-module breakdown
    lines.append("---\n")
    lines.append("## Module-by-Module Reuse Candidates\n")
    lines.append("> Slides are ranked by relevance score (number of matching keywords).")
    lines.append("> Higher score = stronger match to the module topic.\n")

    for mod_name, mod_data in WORKSHOP1_MODULES.items():
        lines.append(f"### {mod_name}\n")

        sorted_slides = sorted(mod_data["slides"], key=lambda s: s["relevance_score"], reverse=True)

        if not sorted_slides:
            lines.append("*No strong matches found in existing decks.*\n")
            continue

        lines.append("| Score | Source Deck | Slide # | Title | Key Matches | Content |")
        lines.append("|-------|-----------|---------|-------|-------------|---------|")

        for slide in sorted_slides[:15]:  # Top 15 per module
            title_clean = slide["title"].replace("|", "\\|")[:80]
            matched = ", ".join(slide["matched_keywords"][:5])
            preview = slide["text_preview"].replace("|", " ").replace("\n", " ")[:120]
            has_extras = []
            if slide["has_table"]:
                has_extras.append("📊")
            if slide["has_image"]:
                has_extras.append("🖼️")
            extras = " ".join(has_extras)
            lines.append(
                f"| **{slide['relevance_score']}** | {slide['file'][:45]} | {slide['slide_num']} "
                f"| {title_clean} | {matched} | {extras} {preview} |"
            )

        lines.append("")

    # High-value slides (score >= 4) across all modules
    lines.append("---\n")
    lines.append("## 🏆 Top Reuse Candidates (Score ≥ 4)\n")
    lines.append("These slides have the strongest overlap with Workshop 1 topics:\n")

    high_value = []
    for mod_name, mod_data in WORKSHOP1_MODULES.items():
        for slide in mod_data["slides"]:
            if slide["relevance_score"] >= 4:
                high_value.append((mod_name, slide))

    high_value.sort(key=lambda x: x[1]["relevance_score"], reverse=True)

    if high_value:
        lines.append("| Score | Module | Source Deck | Slide # | Title |")
        lines.append("|-------|--------|-----------|---------|-------|")
        for mod_name, slide in high_value[:30]:
            title_clean = slide["title"].replace("|", "\\|")[:80]
            lines.append(
                f"| **{slide['relevance_score']}** | {mod_name.split(':')[0]} | "
                f"{slide['file'][:45]} | {slide['slide_num']} | {title_clean} |"
            )
    else:
        lines.append("*No slides scored ≥ 4. See per-module tables above for lower-scoring matches.*\n")

    lines.append("")

    # Recommendation section
    lines.append("---\n")
    lines.append("## 💡 Recommendations\n")
    lines.append("### Reuse Strategy\n")
    lines.append("1. **Direct reuse** (Score ≥ 5): Copy slide as-is, minor text tweaks only")
    lines.append("2. **Adapt** (Score 3-4): Good foundation, needs workshop-specific framing")
    lines.append("3. **Reference only** (Score 2): Topic overlap but likely needs full rebuild\n")
    lines.append("### Next Steps\n")
    lines.append("1. Open the top-scoring PPTX files and review the identified slides visually")
    lines.append("2. Copy candidate slides into a new Workshop 1 deck")
    lines.append("3. Add workshop-specific framing (learning objectives, transitions, exercises)")
    lines.append("4. Fill gaps where no existing slide covers the topic\n")

    # Full slide inventory appendix
    lines.append("---\n")
    lines.append("## 📋 Appendix: Full Slide Inventory\n")
    for fname in file_stats:
        lines.append(f"### {fname}\n")
        file_slides = [s for s in all_slides if s["file"] == fname]
        lines.append("| # | Title | Words | Table | Image |")
        lines.append("|---|-------|-------|-------|-------|")
        for slide in file_slides:
            title_clean = slide["title"].replace("|", "\\|")[:80]
            tbl = "✅" if slide["has_table"] else ""
            img = "✅" if slide["has_image"] else ""
            lines.append(f"| {slide['slide_num']} | {title_clean} | {slide['word_count']} | {tbl} | {img} |")
        lines.append("")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    return output_path


def main():
    repo_root = Path(__file__).parent.parent
    pptx_dir = repo_root / "presentations"
    output_path = repo_root / "scripts" / "output" / "workshop1-slide-reuse-map.md"

    pptx_files = sorted(pptx_dir.glob("*.pptx"))
    if not pptx_files:
        print("No PPTX files found in presentations/")
        sys.exit(1)

    print(f"Found {len(pptx_files)} PPTX files to analyze:\n")

    all_slides = []
    for pptx_path in pptx_files:
        slides = analyze_presentation(str(pptx_path))
        all_slides.extend(slides)

    print(f"\nTotal slides across all decks: {len(all_slides)}")
    print("Mapping to Workshop 1 modules...")

    report_path = generate_report(all_slides, str(output_path))
    print(f"\n✅ Report generated: {report_path}")


if __name__ == "__main__":
    main()
