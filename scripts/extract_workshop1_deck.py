#!/usr/bin/env python3
"""
Extract top-scoring slides from existing presentations into a new Workshop 1 PPTX deck.

Reads the source PPTX files, identifies the best slides for each Workshop 1 module
based on keyword relevance scoring, and copies them into a structured new deck with
section dividers and placeholder slides for gaps.

Usage:
    python scripts/extract_workshop1_deck.py

Output:
    workshops/01-landing-zone-fundamentals/Workshop1-AI-Landing-Zone-Fundamentals.pptx
"""

import copy
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI


# ──────────────────────────────────────────────────────────────────────
# Configuration: source decks and per-module slide selections
# ──────────────────────────────────────────────────────────────────────

REPO_ROOT = Path(__file__).parent.parent

PRESENTATIONS_DIR = REPO_ROOT / "presentations"

OUTPUT_PATH = (
    REPO_ROOT
    / "workshops"
    / "01-landing-zone-fundamentals"
    / "Workshop1-AI-Landing-Zone-Fundamentals.pptx"
)

# Source decks (only the ones python-pptx can open)
SOURCE_DECKS = {
    "external": "AI Landing Zones - External - Presentation.pptx",
    "apim": "AI Landing Zones for APIM - Customer Ready.pptx",
    "foundry100": "AI Landing Zones for Foundry - Customer Ready - L100-200.pptx",
    "foundry300": "AI Landing Zones for Foundry - Customer Ready - L300-400.pptx",
}

# ──────────────────────────────────────────────────────────────────────
# Workshop 1 structure: ordered list of sections and slide selections
#
# Each entry is either:
#   - A "divider" dict: creates a section title slide
#   - A "source" dict: copies a slide from a source deck
#   - A "placeholder" dict: creates a text-only placeholder slide
# ──────────────────────────────────────────────────────────────────────

WORKSHOP1_SLIDES = [
    # ── Section 0: Title & Welcome ──
    {"type": "divider", "title": "AI Landing Zones Fundamentals",
     "subtitle": "Partner Enablement Workshop\nAI Center of Excellence V2 | Q3 FY2026"},

    {"type": "placeholder", "title": "Agenda",
     "body": (
         "Module 1: What is AI Landing Zone?\n"
         "Module 2: Reference Architectures\n"
         "Module 3: Design Checklist Walkthrough\n"
         "Module 4: Deployment Options\n"
         "Module 5: Partner Engagement Scenarios"
     )},

    {"type": "placeholder", "title": "Learning Objectives",
     "body": (
         "• Explain what Azure AI Landing Zones are\n"
         "• Identify key architecture components\n"
         "• Navigate the Design Checklist\n"
         "• Choose appropriate deployment approach\n"
         "• Plan customer engagements"
     )},

    {"type": "placeholder", "title": "Housekeeping",
     "body": (
         "Duration: 2-3 hours\n"
         "Breaks: 10 min after Module 2\n"
         "Questions: Welcome anytime\n"
         "Materials: Links in chat\n"
         "Recording: [Yes/No]"
     )},

    # ── Section 1: What is AI Landing Zone? ──
    {"type": "divider", "title": "Module 1",
     "subtitle": "What is AI Landing Zone?"},

    # Slide 5 — CAF overview
    {"type": "source", "deck": "external", "slide": 5,
     "note": "Microsoft Cloud Adoption Framework for Azure"},
    {"type": "source", "deck": "external", "slide": 6,
     "note": "CAF for Azure (detailed)"},
    {"type": "source", "deck": "external", "slide": 7,
     "note": "Cloud Adoption Framework for AI Adoption"},

    # Slide 36 — AI Landing Zone overview (Score 7)
    {"type": "source", "deck": "external", "slide": 36,
     "note": "AI Landing Zone Overview — top match (Score 7)"},

    # WAF alignment
    {"type": "source", "deck": "external", "slide": 9,
     "note": "Azure Well-Architected Framework"},
    {"type": "source", "deck": "external", "slide": 10,
     "note": "WAF – AI Workload"},

    # Platform vs Application LZ
    {"type": "source", "deck": "external", "slide": 21,
     "note": "Platform LZs and Application LZs"},
    {"type": "source", "deck": "external", "slide": 19,
     "note": "What's a workload-specific landing zone?"},
    {"type": "source", "deck": "external", "slide": 20,
     "note": "Azure landing zone conceptual architecture"},

    # CAF AI scenario positioning
    {"type": "source", "deck": "apim", "slide": 6,
     "note": "Comprehensive AI guidance — CAF + WAF (Score 5)"},

    # Agent types and use cases
    {"type": "source", "deck": "external", "slide": 71,
     "note": "Let's start at the basics: Definitions"},
    {"type": "source", "deck": "external", "slide": 72,
     "note": "Apps, agents and tools — What are your options"},
    {"type": "source", "deck": "external", "slide": 32,
     "note": "Agentic AI System Architecture Logical View"},

    # Agent taxonomy / decision tree (placeholder — likely in encrypted CAF deck)
    {"type": "placeholder", "title": "AI Agent Types (CAF)",
     "body": (
         "Productivity Agents — Retrieve & synthesize information\n"
         "Action Agents — Perform specific tasks in workflows\n"
         "Automation Agents — Multi-step processes, minimal oversight\n\n"
         "Source: CAF AI Agent Adoption\n"
         "https://learn.microsoft.com/azure/cloud-adoption-framework/ai-agents/"
     )},

    {"type": "placeholder", "title": "When to Use Agents vs. Classic RAG",
     "body": (
         "DON'T use agents when:\n"
         "• Task is structured, predictable, rule-based → Deterministic code\n"
         "• Goal is static knowledge retrieval / Q&A → Classic RAG\n\n"
         "USE agents when:\n"
         "• Dynamic decision-making (multi-step reasoning)\n"
         "• Complex orchestration (chaining tools, APIs)\n"
         "• Adaptive behavior (ambiguous inputs)\n\n"
         "Source: CAF AI Agent Business Plan\n"
         "https://learn.microsoft.com/azure/cloud-adoption-framework/ai-agents/"
         "business-strategy-plan"
     )},

    {"type": "placeholder", "title": "Key Takeaways — Module 1",
     "body": (
         "• AI Landing Zone = production-ready AI foundation\n"
         "• Supports both classic RAG and AI agent workloads\n"
         "• Aligns with CAF and WAF\n"
         "• Use agent decision tree to qualify workload type\n"
         "• Accelerates enterprise AI adoption"
     )},

    # ── Section 2: Reference Architectures ──
    {"type": "divider", "title": "Module 2",
     "subtitle": "Reference Architectures"},

    # Two architecture options
    {"type": "source", "deck": "external", "slide": 40,
     "note": "Reference Architectures — with/without platform LZ (Score 4)"},

    # Golden Path architecture
    {"type": "source", "deck": "external", "slide": 33,
     "note": "Recommended 'Golden Path' Architecture (Score 5)"},

    # AI Foundry / Project diagram
    {"type": "source", "deck": "external", "slide": 337,
     "note": "AI Foundry Project — services overview (Score 4)"},
    {"type": "source", "deck": "foundry100", "slide": 25,
     "note": "Azure AI Foundry — platform overview"},

    # Key components — compute
    {"type": "source", "deck": "external", "slide": 48,
     "note": "Build AI Apps & Agents, your way"},
    {"type": "source", "deck": "external", "slide": 49,
     "note": "Choosing the Right Azure Platform for Your AI Workload"},

    # Agent Framework
    {"type": "source", "deck": "external", "slide": 109,
     "note": "Microsoft Agent Framework (Score 4)"},

    # Security & governance in architecture
    {"type": "source", "deck": "external", "slide": 25,
     "note": "Security, compliance, and governance"},

    {"type": "placeholder", "title": "Choosing Your Architecture",
     "body": (
         "Factor              | With Platform LZ | Without Platform LZ\n"
         "Existing ALZ        | ✅ Recommended    | ⚠️ Consider migration\n"
         "Greenfield          | ⚠️ More setup     | ✅ Faster start\n"
         "Enterprise          | ✅ Best fit        | ⚠️ May need upgrade\n"
         "PoC/Pilot           | ⚠️ Overkill       | ✅ Right-sized"
     )},

    {"type": "placeholder", "title": "Key Takeaways — Module 2",
     "body": (
         "• Two validated architecture options\n"
         "• Choose based on customer context\n"
         "• All components are Azure PaaS\n"
         "• Security-first design by default"
     )},

    # ── BREAK ──
    {"type": "divider", "title": "☕ Break",
     "subtitle": "10 minutes — we'll resume with Module 3"},

    # ── Section 3: Design Checklist Walkthrough ──
    {"type": "divider", "title": "Module 3",
     "subtitle": "Design Checklist Walkthrough"},

    # Design areas overview
    {"type": "source", "deck": "external", "slide": 18,
     "note": "Azure Landing Zones — Design areas (Score 5)"},
    {"type": "source", "deck": "external", "slide": 16,
     "note": "Azure landing zones overview (Score 4)"},
    {"type": "source", "deck": "external", "slide": 17,
     "note": "Design Principles"},

    # Customer pain points — why checklist matters
    {"type": "source", "deck": "external", "slide": 35,
     "note": "Customer pain points (Score 5)"},
    {"type": "source", "deck": "external", "slide": 37,
     "note": "Key features (Score 4)"},

    # Design area deep dives
    {"type": "source", "deck": "external", "slide": 47,
     "note": "Compute — Design Considerations, Recommendations, Decisions"},
    {"type": "source", "deck": "external", "slide": 225,
     "note": "Design Area Considerations, Recommendations, Decisions (Score 4)"},
    {"type": "source", "deck": "external", "slide": 263,
     "note": "Agent Security (Score 4)"},
    {"type": "source", "deck": "external", "slide": 294,
     "note": "Design Area — Cost Considerations (Score 4)"},

    # Foundry security / governance
    {"type": "source", "deck": "external", "slide": 57,
     "note": "Microsoft Foundry 3/3 — comprehensive capabilities (Score 5)"},

    {"type": "placeholder", "title": "How to Use the Design Checklist",
     "body": (
         "1. Assessment — Walk through with customer\n"
         "2. Gap Analysis — Identify missing items\n"
         "3. Prioritization — P0/P1/P2 classification\n"
         "4. Implementation — Track completion\n"
         "5. Validation — Post-deployment review\n\n"
         "Full Checklist: github.com/Azure/AI-Landing-Zones/blob/main/"
         "docs/AI-Landing-Zones-Design-Checklist.md"
     )},

    {"type": "placeholder", "title": "Key Takeaways — Module 3",
     "body": (
         "• Design Checklist is your primary assessment tool\n"
         "• 40+ recommendations across 10 design areas\n"
         "• Use with every customer engagement\n"
         "• Reference, don't recreate"
     )},

    # ── Section 4: Deployment Options ──
    {"type": "divider", "title": "Module 4",
     "subtitle": "Deployment Options"},

    # ALZ Customer Journey (Bicep/TF/Portal paths)
    {"type": "source", "deck": "external", "slide": 22,
     "note": "Azure Landing Zone Customer Journey (Score 4)"},

    # AVM modules
    {"type": "source", "deck": "external", "slide": 43,
     "note": "Solution nesting — leveraging AVM pattern modules (Score 4)"},

    # Deployment types / cost
    {"type": "source", "deck": "external", "slide": 160,
     "note": "Choosing the best deployment location"},
    {"type": "source", "deck": "external", "slide": 295,
     "note": "Saving with Azure AI Foundry PTU reservations"},

    {"type": "placeholder", "title": "Four Deployment Paths",
     "body": (
         "Path | Tool         | Speed       | Customization\n"
         "A    | azd up       | ~45 min     | Limited\n"
         "B    | Bicep        | ~30-60 min  | Full\n"
         "C    | Terraform    | ~30-60 min  | Full\n"
         "D    | Portal       | ~30-45 min  | Limited\n\n"
         "See IaC Decision Framework for detailed guidance"
     )},

    {"type": "placeholder", "title": "⚠️ Cost Warning — Know Before You Deploy",
     "body": (
         "Standard deployment: ~$2,128-3,098/month\n"
         "Dev/Test optimized: ~$980-1,350/month\n"
         "Maximum optimization: ~$640-900/month\n\n"
         "High-cost services: Firewall ($400-500), App Gateway ($350-450),\n"
         "Bastion (~$140), AI Search S1 (~$250), Cosmos DB (~$175)\n\n"
         "⚠️ DELETE all resources immediately after workshops!\n"
         "Cost Guide: github.com/Azure/AI-Landing-Zones/blob/main/"
         "docs/AI-Landing-Zones-Cost-Guide.md"
     )},

    {"type": "placeholder", "title": "IaC Decision Framework",
     "body": (
         "Need fastest deployment? → azd up\n"
         "Multi-cloud strategy? → Terraform\n"
         "Azure-native team? → Bicep\n"
         "No-code / quick demo? → Portal (Deploy to Azure)\n\n"
         "See docs/IAC-DECISION-FRAMEWORK.md for full decision tree"
     )},

    {"type": "placeholder", "title": "Key Takeaways — Module 4",
     "body": (
         "• Multiple valid deployment paths\n"
         "• Choose based on customer context\n"
         "• All paths are production-ready\n"
         "• Watch costs — delete resources after labs"
     )},

    # ── Section 5: Partner Engagement Scenarios ──
    {"type": "divider", "title": "Module 5",
     "subtitle": "Partner Engagement Scenarios"},

    # Take the next step (from APIM deck)
    {"type": "source", "deck": "apim", "slide": 108,
     "note": "Take the next step with Azure — engagement paths"},

    # Agent framework for customer scenarios
    {"type": "source", "deck": "apim", "slide": 43,
     "note": "Agent frameworks and agentic architecture (Score 4)"},

    {"type": "placeholder", "title": "Scenario 1 — Enterprise with Existing ALZ",
     "body": (
         "Context: Large org with existing Azure Landing Zones\n"
         "Recommendation: Bicep/Terraform + Platform LZ architecture\n"
         "Focus: Networking, Identity integration\n\n"
         "Key Questions:\n"
         "• What's their current ALZ configuration?\n"
         "• What governance policies exist?\n"
         "• Integration requirements?"
     )},

    {"type": "placeholder", "title": "Scenario 2 — Greenfield / PoC",
     "body": (
         "Context: New to Azure or exploring AI\n"
         "Recommendation: azd up (Deploy-Your-AI-App)\n"
         "Focus: Fast deployment, validation\n\n"
         "Key Questions:\n"
         "• Timeline for PoC?\n"
         "• Expansion plans?\n"
         "• Who needs access?"
     )},

    {"type": "placeholder", "title": "Scenario 3 — Regulated Industry",
     "body": (
         "Context: Healthcare, finance, government\n"
         "Recommendation: Bicep + extensive security review\n"
         "Focus: Security (S-R1-5), Governance (G-R1-5)\n\n"
         "Key Questions:\n"
         "• Compliance requirements?\n"
         "• Data residency needs?\n"
         "• Audit requirements?"
     )},

    {"type": "placeholder", "title": "Scenario 4 — Cost-Sensitive",
     "body": (
         "Context: Limited budget, cost optimization critical\n"
         "Recommendation: Minimal deployment + Cost checklist\n"
         "Focus: Cost (CO-R1-4), Compute (C-R1)\n\n"
         "Key Questions:\n"
         "• Budget constraints?\n"
         "• Workload predictability?\n"
         "• Scale requirements?"
     )},

    {"type": "placeholder", "title": "Scenario 5 — Customer Building AI Agents",
     "body": (
         "Context: Moving beyond chat/RAG to AI agents\n"
         "Recommendation: Standard Foundry setup + CAF AI Agent guidance\n"
         "Focus: Security, Governance, Compute, Monitoring\n\n"
         "Key Questions:\n"
         "• What tasks will the agent perform?\n"
         "• Single agent or multi-agent?\n"
         "• Who owns governance?\n"
         "• SaaS agents sufficient?"
     )},

    {"type": "placeholder", "title": "Key Takeaways — Module 5",
     "body": (
         "• Tailor approach to customer context\n"
         "• Use Design Checklist as conversation guide\n"
         "• Ask discovery questions upfront\n"
         "• Use CAF agent decision tree for agent workloads"
     )},

    # ── Section 6: Wrap-up ──
    {"type": "divider", "title": "Workshop Summary",
     "subtitle": "What We Covered Today"},

    {"type": "placeholder", "title": "Today We Covered",
     "body": (
         "✅ What AI Landing Zones are and why they matter\n"
         "✅ Classic RAG vs. AI Agents — when to use each\n"
         "✅ Two reference architecture options\n"
         "✅ 10 design areas and key recommendations\n"
         "✅ Four deployment paths\n"
         "✅ Five partner engagement scenarios"
     )},

    {"type": "placeholder", "title": "Key Resources",
     "body": (
         "AI Landing Zones Repo: github.com/Azure/AI-Landing-Zones\n"
         "Deploy-Your-AI-App: github.com/microsoft/Deploy-Your-AI-Application-In-Production\n"
         "Design Checklist: AI-Landing-Zones-Design-Checklist.md\n"
         "Cost Guide: AI-Landing-Zones-Cost-Guide.md\n"
         "CAF AI Agent Adoption: learn.microsoft.com/.../ai-agents/\n"
         "IaC Decision Framework: docs/IAC-DECISION-FRAMEWORK.md"
     )},

    {"type": "placeholder", "title": "Next Steps",
     "body": (
         "1. Practice — Deploy a test environment (see cost warnings!)\n"
         "2. ⚠️ Delete resources — Tear down immediately after practice\n"
         "3. Read — Review Design Checklist in detail\n"
         "4. Explore — CAF AI Agent Adoption guidance\n"
         "5. Apply — Use with next customer engagement\n"
         "6. Continue — Attend Workshop 2 (hands-on deployment)"
     )},

    {"type": "placeholder", "title": "Workshop 2 Preview",
     "body": (
         "From RAG to Agents: Deploying Your First Gen AI Workload\n\n"
         "• Hands-on lab: Deploy Landing Zone + RAG chat app\n"
         "• Configure AI Foundry with standard (private) setup\n"
         "• Understand when to graduate from RAG to agents\n"
         "• Explore Microsoft Foundry agent capabilities\n"
         "• Implement monitoring and observability"
     )},

    {"type": "placeholder", "title": "Q&A",
     "body": (
         "Questions?\n"
         "Feedback?\n\n"
         "Thank you!"
     )},
]


# ──────────────────────────────────────────────────────────────────────
# Slide creation helpers
# ──────────────────────────────────────────────────────────────────────

# Microsoft-ish color scheme
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

# Namespace for relationship attributes in OOXML
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Global counter for unique media part names across all copied slides
_img_counter = 0


def _get_blank_layout(prs):
    """Find a blank slide layout (fewest placeholders)."""
    best = prs.slide_layouts[0]
    best_count = len(list(best.placeholders))
    for layout in prs.slide_layouts:
        count = len(list(layout.placeholders))
        if count < best_count:
            best = layout
            best_count = count
        if count == 0:
            return layout
    return best


def copy_slide(source_prs, slide_index, target_prs):
    """
    Copy a slide from source presentation to target presentation.

    Properly handles image/media relationships by:
    1. Creating new Part objects for each media blob in the target package
    2. Adding relationships from the new slide to those parts
    3. Remapping rId references in shape XML to the new rIds
    """
    global _img_counter

    src_slide = source_prs.slides[slide_index]
    src_part = src_slide.part

    # Add blank slide to target
    layout = _get_blank_layout(target_prs)
    new_slide = target_prs.slides.add_slide(layout)
    new_part = new_slide.part

    # Remove any default placeholder shapes from the blank layout
    sp_tree = new_slide.shapes._spTree
    removable_tags = {'sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'}
    for child in list(sp_tree):
        local_tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if local_tag in removable_tags:
            sp_tree.remove(child)

    # ── Phase 1: Copy relationships and build rId mapping ──
    rId_map = {}  # source_rId -> target_rId

    for src_rel in src_part.rels.values():
        rt = src_rel.reltype

        # Skip layout, notes, tags — we use our own
        if any(skip in rt for skip in ['slideLayout', 'notesSlide', 'tags']):
            continue

        # External relationships (hyperlinks, etc.)
        if src_rel.is_external:
            try:
                new_rId = new_part.rels.get_or_add_ext_rel(rt, src_rel.target_ref)
                rId_map[src_rel.rId] = new_rId
            except Exception:
                pass
            continue

        # Internal relationships (images, media, charts, etc.)
        try:
            src_target = src_rel.target_part
            blob = src_target.blob
            ct = src_target.content_type

            # Determine file extension from content type
            ext = ct.split('/')[-1]
            ext = ext.replace('jpeg', 'jpg')
            if ext.startswith('x-'):
                ext = ext[2:]  # x-emf -> emf, x-wmf -> wmf
            if ext in ('vnd.openxmlformats-officedocument.drawing+xml',):
                ext = 'xml'
            if ext in ('octet-stream',):
                ext = 'bin'

            _img_counter += 1
            new_pn = PackURI(f'/ppt/media/ws1_img{_img_counter}.{ext}')

            # Create a new Part in the target package with the image data
            new_target = Part(new_pn, ct, target_prs.part.package, blob)

            # Create relationship from new slide to the new part
            new_rId = new_part.relate_to(new_target, rt)
            rId_map[src_rel.rId] = new_rId
        except Exception:
            # If we can't copy this part, leave rId unmapped
            pass

    # ── Phase 2: Copy shape XML with rId remapping ──
    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)

        # Remap all relationship rId references in the copied XML
        for node in el.iter():
            # Check r:embed, r:link, r:id (namespace-qualified attributes)
            for r_attr in ('embed', 'link', 'id', 'blip'):
                qname = f'{{{NS_R}}}{r_attr}'
                if qname in node.attrib:
                    old_val = node.attrib[qname]
                    if old_val in rId_map:
                        node.attrib[qname] = rId_map[old_val]

            # Also check plain attributes that look like rId references
            for attr_name, attr_val in list(node.attrib.items()):
                if attr_val in rId_map and '{' not in attr_name:
                    node.attrib[attr_name] = rId_map[attr_val]

        sp_tree.append(el)

    # ── Phase 3: Copy background if present ──
    try:
        src_bg = src_slide.background._element
        if len(src_bg) > 0:
            new_bg_el = copy.deepcopy(src_bg)
            # Remap rId refs in background too
            for node in new_bg_el.iter():
                for r_attr in ('embed', 'link', 'id'):
                    qname = f'{{{NS_R}}}{r_attr}'
                    if qname in node.attrib:
                        old_val = node.attrib[qname]
                        if old_val in rId_map:
                            node.attrib[qname] = rId_map[old_val]
            existing_bg = new_slide.background._element
            existing_bg.getparent().replace(existing_bg, new_bg_el)
    except Exception:
        pass

    return new_slide


def add_divider_slide(prs, title, subtitle=""):
    """Add a section divider/title slide with a colored background."""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title text
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle text
    if subtitle:
        top2 = Inches(4.0)
        height2 = Inches(1.5)
        txBox2 = slide.shapes.add_textbox(left, top2, width, height2)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.LEFT

    return slide


def add_placeholder_slide(prs, title, body):
    """Add a content placeholder slide with title and body text."""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Title bar accent line
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(0.08)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Title
    left = Inches(0.7)
    top = Inches(0.4)
    width = Inches(8.5)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.LEFT

    # Body
    left = Inches(0.7)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5.0)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body.split("\n")):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT

    return slide


def add_source_note(slide, note_text):
    """Add a small source note at the bottom of a slide."""
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(9)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"📌 Source: {note_text}"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.RIGHT


# ──────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("Workshop 1 Slide Deck Builder")
    print("=" * 60)

    # Load source presentations
    source_presentations = {}
    for key, filename in SOURCE_DECKS.items():
        filepath = PRESENTATIONS_DIR / filename
        if not filepath.exists():
            print(f"  ⚠️  Source not found: {filename}")
            continue
        try:
            print(f"  Loading: {filename} ...", flush=True)
            source_presentations[key] = Presentation(str(filepath))
            slide_count = len(source_presentations[key].slides)
            print(f"    → {slide_count} slides")
        except Exception as e:
            print(f"  ⚠️  Failed to load {filename}: {e}")

    if not source_presentations:
        print("No source presentations loaded. Exiting.")
        sys.exit(1)

    # Get slide dimensions from the primary source deck
    primary = source_presentations.get("external")
    if primary:
        slide_width = primary.slide_width
        slide_height = primary.slide_height
    else:
        slide_width = Inches(13.333)  # Widescreen default
        slide_height = Inches(7.5)

    # Create new presentation
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    print(f"\nSlide dimensions: {slide_width/914400:.1f}\" x {slide_height/914400:.1f}\"")
    print(f"\nBuilding Workshop 1 deck ({len(WORKSHOP1_SLIDES)} entries)...\n")

    stats = {"dividers": 0, "placeholders": 0, "copied": 0, "failed": 0}

    for i, entry in enumerate(WORKSHOP1_SLIDES):
        entry_type = entry["type"]

        if entry_type == "divider":
            add_divider_slide(prs, entry["title"], entry.get("subtitle", ""))
            stats["dividers"] += 1
            print(f"  [{i+1:2d}] 📑 Divider: {entry['title']}")

        elif entry_type == "placeholder":
            add_placeholder_slide(prs, entry["title"], entry["body"])
            stats["placeholders"] += 1
            print(f"  [{i+1:2d}] 📝 Placeholder: {entry['title']}")

        elif entry_type == "source":
            deck_key = entry["deck"]
            slide_num = entry["slide"]  # 1-based
            note = entry.get("note", "")

            if deck_key not in source_presentations:
                print(f"  [{i+1:2d}] ⚠️  Source deck '{deck_key}' not loaded — placeholder instead")
                add_placeholder_slide(prs, f"[MISSING] {note}",
                                      f"Source: {SOURCE_DECKS.get(deck_key, deck_key)}\n"
                                      f"Slide: {slide_num}\n"
                                      f"Please copy this slide manually.")
                stats["failed"] += 1
                continue

            src_prs = source_presentations[deck_key]
            if slide_num < 1 or slide_num > len(src_prs.slides):
                print(f"  [{i+1:2d}] ⚠️  Slide {slide_num} out of range for '{deck_key}' — skipping")
                add_placeholder_slide(prs, f"[OUT OF RANGE] {note}",
                                      f"Source: {SOURCE_DECKS.get(deck_key, deck_key)}\n"
                                      f"Slide {slide_num} (deck has {len(src_prs.slides)} slides)")
                stats["failed"] += 1
                continue

            try:
                new_slide = copy_slide(src_prs, slide_num - 1, prs)
                # Add source annotation
                deck_name = SOURCE_DECKS[deck_key][:50]
                add_source_note(new_slide, f"{deck_name} — Slide {slide_num}")
                stats["copied"] += 1
                print(f"  [{i+1:2d}] ✅ Copied: {deck_key} #{slide_num} — {note}")
            except Exception as e:
                print(f"  [{i+1:2d}] ⚠️  Failed to copy {deck_key} #{slide_num}: {e}")
                add_placeholder_slide(prs, f"[COPY FAILED] {note}",
                                      f"Source: {SOURCE_DECKS.get(deck_key, deck_key)}\n"
                                      f"Slide: {slide_num}\n"
                                      f"Error: {e}\n"
                                      f"Please copy this slide manually.")
                stats["failed"] += 1

    # Save
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(str(OUTPUT_PATH))

    total_slides = len(prs.slides)
    print(f"\n{'=' * 60}")
    print(f"✅ Workshop 1 deck saved: {OUTPUT_PATH}")
    print(f"   Total slides: {total_slides}")
    print(f"   Copied from source: {stats['copied']}")
    print(f"   Section dividers: {stats['dividers']}")
    print(f"   Placeholders (to fill): {stats['placeholders']}")
    if stats["failed"]:
        print(f"   ⚠️  Failed copies: {stats['failed']}")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    main()
